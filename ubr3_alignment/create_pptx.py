"""
Create a PowerPoint presentation with all UBR3 analysis figures,
one figure per slide with detailed explanations beside each figure.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

EXPLANATIONS = {
    "ubr3_animal_domain_architecture.png": (
        "UBR3 Domain Architecture Across the Animal Kingdom",
        [
            "This figure displays the full-length UBR3 protein for each of the 13 species "
            "analyzed, ordered from vertebrates (top) to cnidarians (bottom), reflecting "
            "evolutionary distance from humans.",

            "Three conserved functional domains are highlighted on each protein bar:\n"
            "  \u2022 UBR box (purple, aa 118-189): the substrate-recognition domain that "
            "binds N-degron substrates\n"
            "  \u2022 Conserved region 400-600 (blue): a region of unknown function that is "
            "highly conserved across all species\n"
            "  \u2022 RING domain (orange, aa 1232-1376): the E3 ubiquitin ligase catalytic "
            "domain responsible for ubiquitin transfer",

            "The domain positions are mapped from the human reference sequence onto each "
            "ortholog using the multiple sequence alignment. All 13 species retain all three "
            "domains, demonstrating that the core architecture of UBR3 has been preserved "
            "for over 700 million years of animal evolution.",

            "Animal icons illustrate each species, color-coded by taxonomic group. Protein "
            "lengths (in amino acids) are shown on the right.",
        ]
    ),

    "ubr3_combined_conservation_architecture.png": (
        "Conservation Score and Domain Architecture (Aligned Axes)",
        [
            "This figure combines the per-position conservation plot (Panel A) with the "
            "domain architecture across all 13 species (Panel B), sharing the same x-axis "
            "(human amino acid numbering).",

            "Panel A shows the conservation score (fraction of species sharing the most "
            "common residue) smoothed over a 20-residue window. The red dashed line marks "
            "90% conservation and the orange dashed line marks 70%. Colored vertical bands "
            "highlight the three conserved domains.",

            "Panel B shows the full-length UBR3 protein for each species with colored "
            "domain blocks. Because both panels share the same x-axis, you can trace "
            "straight down from each conservation peak to see the corresponding domain "
            "in every species.",

            "This view makes it immediately clear that the conservation peaks correspond "
            "precisely to the functional domains: the UBR box (~118-189), the 400-600 "
            "region, and the RING domain (~1232-1376). The x-axis uses human residue "
            "numbering, which corrects for alignment gaps and gives biologically meaningful "
            "coordinates.",
        ]
    ),

    "ubr3_pub_combined_panels.png": (
        "Combined Publication Figure",
        [
            "This composite figure combines all key panels into a single publication-ready "
            "image.",

            "Panel A shows ClustalW-style multiple sequence alignments for the three most "
            "conserved functional regions: the UBR box domain (aa 118-189), the conserved "
            "region (aa 400-600), and the RING domain (aa 1232-1376).",

            "Panel B shows the domain architecture diagram with colored blocks indicating "
            "domain positions across all 13 species.",

            "Amino acids are colored by biochemical property: acidic (red), basic (blue), "
            "polar (green), hydrophobic (orange), aromatic (brown), and special residues "
            "(purple/gray). Fully conserved positions are marked with '*', strongly similar "
            "with ':', and weakly similar with '.'.",
        ]
    ),

    "ubr3_pub_combined_panels_slide.png": (
        "Combined Publication Figure (slide layout)",
        [
            "Same content as the stacked publication composite, but rearranged for slides: "
            "panels a and b (UBR box and conserved region alignments) appear side by side, "
            "then panel c (RING alignment) full width, then panel d (domain architecture) "
            "full width. The image is built at high width so it stays legible when scaled "
            "to fit a widescreen slide.",

            "Panels a–c show ClustalW-style multiple sequence alignments for the three most "
            "conserved functional regions. Panel d shows domain positions across all 13 species.",

            "Amino acids are colored by biochemical property. Conservation symbols: "
            "'*' identical, ':' strongly similar, '.' weakly similar.",

            "For PowerPoint, prefer the three separate files "
            "ubr3_pub_slide_panel_ab.png, _c.png, and _d.png (one slide each) if the "
            "all-in-one file is still too tall when fitted to a 16×9 slide.",
        ]
    ),

    "ubr3_pub_slide_panel_ab.png": (
        "UBR publication alignments — panels a and b (slide)",
        [
            "Side-by-side ClustalW-style alignments: UBR box (a) and conserved region 400–600 (b), "
            "scaled for a single widescreen slide."
        ]
    ),

    "ubr3_pub_slide_panel_c.png": (
        "UBR publication alignment — panel c — RING domain (slide)",
        [
            "RING domain alignment block (scaled for one widescreen slide)."
        ]
    ),

    "ubr3_pub_slide_panel_d.png": (
        "UBR domain architecture — panel d (slide)",
        [
            "Domain architecture across species (scaled to match panels a–c width)."
        ]
    ),

    "ubr3_conservation_line.png": (
        "Conservation Score Along UBR3 Alignment",
        [
            "This line plot shows the per-position conservation score across the entire "
            "UBR3 alignment (human residue positions on x-axis).",

            "The conservation score at each column is calculated as the fraction of species "
            "sharing the most common amino acid at that position (ignoring gaps). A score "
            "of 1.0 means all species have the identical residue; lower values indicate "
            "divergence.",

            "A sliding window average (typically 10-20 residues) smooths the signal to "
            "reveal broad conserved regions. Prominent peaks correspond to the UBR box "
            "(~120-190), the 400-600 region, and the RING domain (~1232-1376).",

            "This plot was generated using Biopython to parse the Clustal Omega alignment "
            "and Matplotlib for visualization. Regions above the conservation threshold "
            "are highlighted, identifying candidate functional domains.",
        ]
    ),

    "ubr3_conservation_heatmap.png": (
        "Conservation Heatmap",
        [
            "This heatmap visualizes the alignment as a 2D matrix where each row is a "
            "species and each column is an alignment position.",

            "Colors represent amino acid conservation: darker/warmer colors indicate "
            "positions where the residue matches the consensus, while lighter/cooler colors "
            "indicate divergent positions. Gaps appear as a distinct color.",

            "Vertical bands of consistent color reveal regions conserved across all species. "
            "The UBR box, 400-600 region, and RING domain stand out as prominent dark bands.",

            "This view allows rapid visual identification of both globally conserved regions "
            "and species-specific insertions or deletions.",
        ]
    ),

    "ubr3_conserved_only_heatmap.png": (
        "Conserved Regions Only \u2013 Heatmap",
        [
            "This heatmap is filtered to show only the alignment columns that exceed a "
            "70% conservation threshold, removing variable regions.",

            "By excluding poorly conserved positions, this view concentrates on the "
            "functionally important residues and makes it easier to see which species "
            "diverge at specific conserved sites.",

            "The remaining columns cluster into the known functional domains and reveal "
            "additional short conserved motifs scattered throughout the protein.",
        ]
    ),

    "ubr3_conserved_regions_map.png": (
        "Conserved Regions Map",
        [
            "This figure maps all identified conserved regions onto the UBR3 protein "
            "length, showing their positions and extent.",

            "Each conserved region is defined as a stretch of consecutive alignment columns "
            "where the conservation score exceeds 70% for at least a minimum number of "
            "consecutive positions.",

            "The map reveals that conservation is not uniformly distributed: it clusters "
            "heavily in the N-terminal UBR box, the central 400-600 region, and the "
            "C-terminal RING domain, with scattered smaller motifs in between.",
        ]
    ),

    "ubr3_top_regions_annotated.png": (
        "Top Conserved Regions with Amino Acid Annotations",
        [
            "This bar chart ranks the most conserved regions by their average conservation "
            "score, with the actual amino acid consensus sequence annotated.",

            "For each top region, the consensus amino acid sequence is shown, revealing "
            "the specific residues under strongest evolutionary constraint.",

            "Residues that are 100% conserved across all 13 species are likely essential "
            "for protein function \u2013 mutations at these positions would be expected to "
            "disrupt UBR3 activity.",
        ]
    ),

    "ubr3_pub_domain_architecture.png": (
        "Domain Architecture Diagram",
        [
            "This schematic shows the linear domain architecture of UBR3 across all 13 "
            "species, with species names color-coded by taxonomic group.",

            "The protein backbone is shown as a gray bar proportional to the sequence "
            "length. Colored blocks mark the three characterized domains: UBR box (purple), "
            "conserved region 400-600 (blue), and RING domain (orange).",

            "Domain boundaries in each ortholog are determined by mapping the human domain "
            "coordinates through the multiple sequence alignment to find the corresponding "
            "positions in each species.",

            "The consistent domain organization across all species \u2013 from human to sea "
            "anemone \u2013 indicates strong purifying selection on UBR3 structure.",
        ]
    ),

    "ubr3_domain_conservation.png": (
        "Domain-Specific Conservation",
        [
            "This figure compares the average conservation score within each functional "
            "domain versus the overall protein average.",

            "The UBR box and RING domains typically show significantly higher conservation "
            "than the protein average, consistent with their critical enzymatic functions.",

            "The 400-600 region also shows elevated conservation, suggesting it plays an "
            "important but still uncharacterized role in UBR3 function \u2013 possibly "
            "mediating protein-protein interactions or structural stability.",

            "Inter-domain linker regions show the lowest conservation, consistent with "
            "their role as flexible spacers rather than functional elements.",
        ]
    ),

    "ubr3_identity_matrix.png": (
        "Pairwise Sequence Identity Matrix",
        [
            "This matrix shows the percent sequence identity between every pair of species "
            "in the alignment, displayed as a color-coded heatmap.",

            "Closely related species (e.g., Human vs. Mouse) share high identity (typically "
            ">80%), while distant comparisons (e.g., Human vs. Hydra) show lower identity "
            "(typically 25-40%).",

            "The fact that even the most distant species pairs retain >20% identity across "
            "a ~1900 amino acid protein is remarkable and indicates strong functional "
            "constraint on UBR3 throughout animal evolution.",

            "The matrix was computed from the Clustal Omega multiple sequence alignment by "
            "counting identical residue pairs at each aligned position (excluding gaps).",
        ]
    ),

    "ubr3_pub_ubr_box.png": (
        "UBR Box Domain Alignment (aa 118-189)",
        [
            "ClustalW-style alignment of the UBR box domain (human residues 118-189) "
            "across all 13 species.",

            "The UBR box is the substrate-recognition domain of UBR-family E3 ligases. "
            "It binds N-terminal destabilizing residues (N-degrons) of target proteins, "
            "marking them for proteasomal degradation via the N-end rule pathway.",

            "Residues with identical amino acids across all species are highlighted with "
            "colored backgrounds grouped by biochemical property. Conservation symbols: "
            "'*' = identical in all species, ':' = strongly similar (conserved biochemical "
            "property), '.' = weakly similar.",

            "The high density of fully conserved positions (marked with '*') indicates "
            "that the substrate-binding interface has been under very strong purifying "
            "selection. Key zinc-coordinating cysteines and histidines are 100% conserved.",

            "Red triangles mark substrate recognition residues identified from the UBR3-UNC13B "
            "cryo-EM structure: Arg132, Pro140, Cys141, Ala169, Asp171, and Val177. These "
            "residues form the Pro1-binding pocket and contact the N-degron peptide.",
        ]
    ),

    "ubr3_pub_region_400_600.png": (
        "Conserved Region Alignment (aa 400-600)",
        [
            "ClustalW-style alignment of the conserved region spanning human residues "
            "400-600 across all 13 species.",

            "This ~200 amino acid region shows unexpectedly high conservation despite "
            "having no currently characterized function in the literature. Its conservation "
            "level rivals that of the known functional domains.",

            "The alignment reveals numerous fully conserved positions scattered throughout "
            "this region, suggesting it plays an important structural or functional role "
            "that has not yet been experimentally characterized.",

            "This region may represent a novel functional domain, a protein-protein "
            "interaction surface, or a structural scaffold essential for UBR3 folding. "
            "It represents a promising target for future mutagenesis studies.",

            "Red triangles mark key substrate recognition residues from the UBR3-UNC13B "
            "structure: Val444, Val448, Gln449 (Pro1-binding pocket), Trp511 (Leu3/Arg5 "
            "interactions), and Glu577 (Arg5/Lys7 electrostatic contacts). These explain "
            "why this region is so highly conserved.",
        ]
    ),

    "ubr3_pub_ring_domain.png": (
        "RING Domain Alignment (aa 1232-1376)",
        [
            "ClustalW-style alignment of the RING finger domain (human residues 1232-1376) "
            "across all 13 species.",

            "The RING domain is the catalytic core of UBR3's E3 ubiquitin ligase activity. "
            "It coordinates two zinc ions through a conserved pattern of cysteine and "
            "histidine residues (C3HC4 or C3H2C3 type) and directly mediates the transfer "
            "of ubiquitin from E2 conjugating enzymes to substrate proteins.",

            "The zinc-coordinating residues (cysteines and histidines) are 100% conserved "
            "across all 13 species, reflecting their absolute requirement for RING domain "
            "structure and function.",

            "Even non-zinc-coordinating positions show high conservation, indicating that "
            "the overall fold and E2-interaction surface are under strong selection.",

            "Blue markers highlight the four coordinating cysteines (Cys1239, Cys1334, "
            "Cys1360, Cys1363) emphasized in structural studies.",
        ]
    ),

    "ubr3_pub_ubr_box_ppt.png": (
        "UBR Box Domain Alignment — slide version (large type)",
        [
            "Use this file in PowerPoint: same alignment as the standard export, but with "
            "larger amino-acid and species labels, wider columns, higher resolution, and "
            "more prominent conservation symbols (*, :, .).",

            "ClustalW-style alignment of the UBR box domain (human residues 118-189) "
            "across all 13 species.",

            "The UBR box is the substrate-recognition domain of UBR-family E3 ligases. "
            "It binds N-terminal destabilizing residues (N-degrons) of target proteins, "
            "marking them for proteasomal degradation via the N-end rule pathway.",

            "Red triangles mark substrate recognition residues from the UBR3-UNC13B "
            "structure: Arg132, Pro140, Cys141, Ala169, Asp171, and Val177.",
        ]
    ),

    "ubr3_pub_region_400_600_ppt.png": (
        "Conserved Region Alignment — slide version (large type)",
        [
            "Use this file in PowerPoint: enlarged fonts and conservation row for readability "
            "on a projector.",

            "Same region as the standard figure: human residues 400–600 across all 13 species.",

            "Red triangles mark Val444, Val448, Gln449, Trp511, and Glu577 from the "
            "UBR3-UNC13B structure.",
        ]
    ),

    "ubr3_pub_ring_domain_ppt.png": (
        "RING Domain Alignment — slide version (large type)",
        [
            "Use this file in PowerPoint: enlarged fonts for RING residues and clearer "
            "conservation markers (* larger and bold).",

            "RING domain alignment (human residues 1232–1376). Blue markers: Cys1239, "
            "Cys1334, Cys1360, Cys1363.",
        ]
    ),
}

BLOCK_EXPLANATION = (
    "Annotated Alignment Block",
    [
        "This panel shows a portion of the full multiple sequence alignment of UBR3 "
        "across 13 species, displayed in blocks of ~60 residues per line.",

        "Each row represents one species. Amino acids are colored by conservation: "
        "positions identical across all species are highlighted with colored backgrounds "
        "matching their biochemical group. Gaps (insertions/deletions) are shown as '-'.",

        "Human residue numbers are marked at the top of each block for reference. "
        "Conservation symbols below each block indicate: '*' = fully conserved, "
        "':' = strongly similar, '.' = weakly similar.",

        "These blocks collectively show the full-length alignment, allowing detailed "
        "inspection of any region of interest.",
    ]
)

REGION_EXPLANATION = (
    "Individual Conserved Region",
    [
        "This panel shows a zoomed-in view of a specific conserved region identified "
        "by the conservation analysis (regions exceeding 70% average conservation score).",

        "The alignment is displayed for only the columns within this region, making it "
        "easy to inspect which specific amino acids are conserved and which species "
        "diverge at each position.",

        "These small conserved motifs may represent functional elements such as "
        "phosphorylation sites, protein-protein interaction motifs, post-translational "
        "modification sites, or structural elements critical for protein folding.",

        "The human residue positions are indicated in the title, allowing cross-reference "
        "with protein databases (UniProt, PDB) for further functional annotation.",
    ]
)


SLIDE_GROUPS = [
    ("Overview", [
        "ubr3_pub_slide_panel_ab.png",
        "ubr3_pub_slide_panel_c.png",
        "ubr3_pub_slide_panel_d.png",
        "ubr3_combined_conservation_architecture.png",
        "ubr3_animal_domain_architecture.png",
        "ubr3_pub_combined_panels.png",
        "ubr3_pub_combined_panels_slide.png",
    ]),
    ("Conservation Analysis", [
        "ubr3_conservation_line.png",
        "ubr3_conservation_heatmap.png",
        "ubr3_conserved_only_heatmap.png",
        "ubr3_conserved_regions_map.png",
        "ubr3_top_regions_annotated.png",
    ]),
    ("Domain Analysis", [
        "ubr3_pub_domain_architecture.png",
        "ubr3_domain_conservation.png",
        "ubr3_identity_matrix.png",
    ]),
    ("Publication Alignments", [
        "ubr3_pub_ubr_box.png",
        "ubr3_pub_region_400_600.png",
        "ubr3_pub_ring_domain.png",
    ]),
    ("Full Alignment Blocks", [
        f"ubr3_annotated_block_{i:02d}.png" for i in range(1, 13)
    ]),
    ("Individual Conserved Regions", [
        "ubr3_region_01_h120-149.png",
        "ubr3_region_02_h156-174.png",
        "ubr3_region_03_h393-409.png",
        "ubr3_region_04_h415-430.png",
        "ubr3_region_05_h435-458.png",
        "ubr3_region_06_h469-483.png",
        "ubr3_region_07_h504-527.png",
        "ubr3_region_08_h547-587.png",
        "ubr3_region_09_h626-637.png",
        "ubr3_region_10_h684-699.png",
        "ubr3_region_11_h975-974.png",
        "ubr3_region_12_h1015-1014.png",
        "ubr3_region_13_h1048-1047.png",
        "ubr3_region_14_h1050-1049.png",
        "ubr3_region_15_h1065-1064.png",
        "ubr3_region_16_h1119-1118.png",
        "ubr3_region_17_h1252-1264.png",
        "ubr3_region_18_h1590-1589.png",
        "ubr3_region_19_h1615-1615.png",
        "ubr3_region_20_h1623-1622.png",
    ]),
]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.25)


def add_section_slide(prs, section_title, section_desc):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(2.0), SLIDE_W - Inches(2), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    if section_desc:
        p2 = tf.add_paragraph()
        p2.text = section_desc
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(16)


def add_figure_slide(prs, img_path, title, explanation_paragraphs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_h = Inches(0.65)
    txTitle = slide.shapes.add_textbox(
        MARGIN, Inches(0.08), SLIDE_W - 2 * MARGIN, title_h)
    tf = txTitle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    p.alignment = PP_ALIGN.LEFT

    with Image.open(img_path) as img:
        img_w, img_h = img.size

    text_panel_w = Inches(4.0)
    img_area_w = SLIDE_W - text_panel_w - 3 * MARGIN
    img_area_h = SLIDE_H - title_h - Inches(0.3)

    scale_w = img_area_w / img_w
    scale_h = img_area_h / img_h
    scale = min(scale_w, scale_h)

    final_w = int(img_w * scale)
    final_h = int(img_h * scale)

    img_left = MARGIN
    img_top = int(title_h + (img_area_h - final_h) / 2)

    slide.shapes.add_picture(img_path, img_left, img_top, final_w, final_h)

    text_left = SLIDE_W - text_panel_w - MARGIN
    text_top = title_h + Inches(0.1)
    text_h = SLIDE_H - title_h - Inches(0.4)

    txBox = slide.shapes.add_textbox(text_left, text_top, text_panel_w, text_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    for idx, para_text in enumerate(explanation_paragraphs):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para_text
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(8) if idx > 0 else Pt(0)
        p.space_after = Pt(4)
        p.alignment = PP_ALIGN.LEFT


SECTION_DESCRIPTIONS = {
    "Overview": "Key summary figures of UBR3 conservation across 13 animal species",
    "Conservation Analysis": "Quantitative analysis of sequence conservation along the UBR3 protein",
    "Domain Analysis": "Functional domain comparison and pairwise identity",
    "Publication Alignments": "ClustalW-style alignments of the three major conserved domains",
    "Full Alignment Blocks": "Complete annotated multiple sequence alignment in sequential blocks",
    "Individual Conserved Regions": "Zoomed-in views of each conserved region identified by the analysis",
}


def get_explanation(fname):
    if fname in EXPLANATIONS:
        return EXPLANATIONS[fname]
    if fname.startswith("ubr3_annotated_block_"):
        block_num = fname.replace("ubr3_annotated_block_", "").replace(".png", "")
        title = f"Annotated Alignment Block {int(block_num)}"
        return (title, BLOCK_EXPLANATION[1])
    if fname.startswith("ubr3_region_"):
        parts = fname.replace("ubr3_region_", "").replace(".png", "")
        num_part = parts.split("_")[0]
        pos_part = parts.split("_h")[1] if "_h" in parts else parts
        title = f"Conserved Region {int(num_part)}: Human Position {pos_part}"
        return (title, REGION_EXPLANATION[1])
    return (fname.replace(".png", "").replace("_", " ").title(), [])


def main():
    print("=" * 60)
    print("Creating UBR3 Analysis PowerPoint (with explanations)")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = RGBColor(0x0D, 0x47, 0xA1)

    txBox = title_slide.shapes.add_textbox(
        Inches(1.5), Inches(1.5), SLIDE_W - Inches(3), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "UBR3 Sequence Conservation Analysis"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Cross-species multiple sequence alignment and domain architecture"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "13 species spanning Mammalia to Cnidaria (~700 million years of evolution)"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(10)

    p4 = tf.add_paragraph()
    p4.text = (
        "Species panel: H. sapiens, M. musculus, X. laevis, P. marinus, "
        "C. intestinalis, B. floridae, S. purpuratus, O. sinensis, "
        "D. melanogaster, A. mellifera, P. tepidariorum, H. vulgaris, N. vectensis"
    )
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(20)

    p5 = tf.add_paragraph()
    p5.text = (
        "Methods: Sequences retrieved from NCBI (RefSeq), aligned with Clustal Omega (EBI), "
        "analyzed and visualized with Python (Biopython, Matplotlib)"
    )
    p5.font.size = Pt(11)
    p5.font.color.rgb = RGBColor(0x64, 0xB5, 0xF6)
    p5.alignment = PP_ALIGN.CENTER
    p5.space_before = Pt(10)

    slide_count = 1
    for section_title, figures in SLIDE_GROUPS:
        desc = SECTION_DESCRIPTIONS.get(section_title, "")
        add_section_slide(prs, section_title, desc)
        slide_count += 1
        print(f"\n  Section: {section_title}")

        for fname in figures:
            img_path = os.path.join(OUTPUT_DIR, fname)
            if not os.path.exists(img_path):
                print(f"    SKIP (not found): {fname}")
                continue
            title, paragraphs = get_explanation(fname)
            add_figure_slide(prs, img_path, title, paragraphs)
            slide_count += 1
            print(f"    Added: {fname}")

    out_path = os.path.join(OUTPUT_DIR, "UBR3_Analysis_Figures.pptx")
    prs.save(out_path)
    print(f"\n{'='*60}")
    print(f"Saved: {out_path}")
    print(f"Total slides: {slide_count}")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()

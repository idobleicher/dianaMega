"""
Create a PowerPoint presentation with all RPL22 / RPL22L1 analysis figures,
one figure per slide with detailed explanations.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

EXPLANATIONS = {
    # ── RPL22 figures ─────────────────────────────────────────────────────
    "rpl22_animal_domain_architecture.png": (
        "RPL22 Domain Architecture Across the Animal Kingdom",
        [
            "This figure displays the full-length RPL22 (eL22) ribosomal protein for each "
            "of the species analyzed, ordered from vertebrates (top) to cnidarians (bottom).",

            "The L22e domain (purple) is the core ribosomal protein fold that mediates "
            "RNA binding and incorporation into the 60S ribosomal subunit. It spans most "
            "of the protein (approximately residues 15-128 in human).",

            "RPL22 is a highly conserved component of the eukaryotic 60S ribosomal subunit. "
            "Beyond its ribosomal role, it binds EBV-encoded small RNAs (EBERs) and has "
            "reported roles in cell cycle regulation and hematopoiesis.",

            "Animal icons illustrate each species, color-coded by taxonomic group. Protein "
            "lengths (in amino acids) are shown on the right.",
        ]
    ),

    "rpl22l1_animal_domain_architecture.png": (
        "RPL22L1 Domain Architecture Across the Animal Kingdom",
        [
            "This figure displays the full-length RPL22L1 (eL22-like 1) protein for each "
            "species, ordered from vertebrates to cnidarians.",

            "Like RPL22, RPL22L1 contains a single L22e domain (purple, approximately "
            "residues 16-122 in human) that constitutes nearly the entire protein.",

            "RPL22L1 is a paralog of RPL22 that can substitute for RPL22 in the ribosome. "
            "RPL22 directly represses RPL22L1 expression by binding an internal stem-loop "
            "in the RPL22L1 mRNA, creating a regulatory feedback mechanism.",

            "When RPL22 is lost (as frequently occurs in microsatellite-instable cancers), "
            "RPL22L1 is upregulated and incorporated into ribosomes, maintaining translational "
            "capacity.",
        ]
    ),

    "rpl22_combined_conservation_architecture.png": (
        "RPL22 Conservation Score and Domain Architecture",
        [
            "This figure combines the per-position conservation plot (Panel A) with the "
            "domain architecture across species (Panel B), sharing the same x-axis.",

            "Panel A shows the conservation score (fraction of species sharing the most "
            "common residue) smoothed over a window. The L22e domain region is highlighted "
            "in purple. Red and orange dashed lines mark 90% and 70% conservation.",

            "Panel B shows the full-length RPL22 protein for each species with the L22e "
            "domain colored. Because both panels share the same x-axis, conservation peaks "
            "can be traced directly to domain positions.",

            "The high conservation throughout the L22e domain reflects its essential role "
            "in ribosome assembly and function, while the short N-terminal extension shows "
            "more variability.",
        ]
    ),

    "rpl22l1_combined_conservation_architecture.png": (
        "RPL22L1 Conservation Score and Domain Architecture",
        [
            "Combined conservation plot (Panel A) and domain architecture (Panel B) for "
            "RPL22L1, sharing the same x-axis (human residue numbering).",

            "RPL22L1 shows a similar conservation pattern to RPL22 within the L22e domain, "
            "but with some notable differences in the N-terminal region that may relate to "
            "its distinct regulatory properties.",

            "The L22e domain is well-conserved across all species, consistent with RPL22L1's "
            "ability to functionally substitute for RPL22 in the ribosome.",
        ]
    ),

    "rpl22_pub_full_alignment.png": (
        "RPL22 Full Protein Alignment",
        [
            "ClustalW-style alignment of the complete RPL22 protein across all species.",

            "Residues with identical amino acids across all species are highlighted with "
            "colored backgrounds grouped by biochemical property: acidic (red), basic (blue), "
            "polar (green), hydrophobic (orange), aromatic (brown), special (purple/gray).",

            "Conservation symbols: '*' = identical in all species, ':' = strongly similar "
            "(conserved biochemical property), '.' = weakly similar.",

            "The high density of '*' symbols throughout the alignment demonstrates that "
            "RPL22 is one of the most conserved ribosomal proteins, reflecting its essential "
            "role in the 60S subunit.",
        ]
    ),

    "rpl22l1_pub_full_alignment.png": (
        "RPL22L1 Full Protein Alignment",
        [
            "ClustalW-style alignment of the complete RPL22L1 protein across all species.",

            "The alignment reveals which positions are conserved across the RPL22L1 paralogs "
            "and which have diverged. Comparison with the RPL22 alignment can identify "
            "paralog-specific residues.",

            "Conservation symbols follow the same convention: '*' = identical, ':' = strongly "
            "similar, '.' = weakly similar.",
        ]
    ),

    "rpl22_pub_l22e_domain.png": (
        "RPL22 L22e Domain Alignment",
        [
            "ClustalW-style alignment focused on the L22e domain of RPL22, which constitutes "
            "the core ribosomal protein fold.",

            "The L22e domain mediates binding to 28S rRNA and is required for stable "
            "incorporation into the 60S ribosomal subunit. It adopts a characteristic "
            "alpha/beta fold conserved across eukaryotes.",

            "The extremely high conservation within this domain reflects strong purifying "
            "selection on the ribosome-binding interface.",
        ]
    ),

    "rpl22l1_pub_l22e_domain.png": (
        "RPL22L1 L22e Domain Alignment",
        [
            "ClustalW-style alignment of the L22e domain in RPL22L1.",

            "This domain enables RPL22L1 to substitute for RPL22 in the ribosome. "
            "Comparing this alignment with the RPL22 L22e domain alignment reveals "
            "which positions are conserved in both paralogs (essential for ribosome "
            "function) versus those that differ (potentially related to paralog-specific "
            "regulation).",
        ]
    ),

    "rpl22_pub_domain_architecture.png": (
        "RPL22 Domain Architecture Diagram",
        [
            "Schematic showing the linear domain architecture of RPL22 across all species, "
            "with species names color-coded by taxonomic group.",

            "The protein backbone (gray bar) is proportional to sequence length. The L22e "
            "domain (purple) covers most of the protein in all species.",

            "Domain boundaries are mapped from human coordinates through the multiple "
            "sequence alignment to find corresponding positions in each ortholog.",
        ]
    ),

    "rpl22l1_pub_domain_architecture.png": (
        "RPL22L1 Domain Architecture Diagram",
        [
            "Linear domain architecture of RPL22L1 across all species.",

            "Like RPL22, the L22e domain dominates the protein architecture. The consistent "
            "domain organization across species demonstrates that the paralog duplication "
            "preserved the core ribosomal protein structure.",
        ]
    ),

    "rpl22_pub_combined_panels.png": (
        "RPL22 Combined Publication Figure",
        [
            "Composite figure combining ClustalW-style alignments (Panel A) with the "
            "domain architecture diagram (Panel B) for RPL22.",

            "Panel A shows the full protein and L22e domain alignments. Panel B shows "
            "domain positions across all species.",
        ]
    ),

    "rpl22l1_pub_combined_panels.png": (
        "RPL22L1 Combined Publication Figure",
        [
            "Composite figure combining ClustalW-style alignments (Panel A) with the "
            "domain architecture diagram (Panel B) for RPL22L1.",
        ]
    ),

    # ── Conservation analysis figures ─────────────────────────────────────
    "rpl22_conservation_line.png": (
        "RPL22 Conservation Score Along Alignment",
        [
            "Per-position conservation score across the RPL22 alignment (human residue "
            "positions on x-axis).",

            "The conservation score at each column is the fraction of species sharing the "
            "most common amino acid (ignoring gaps). A sliding window smooths the signal.",

            "RPL22 shows high conservation throughout most of its length, with the L22e "
            "domain region showing particularly strong conservation (>90%).",
        ]
    ),

    "rpl22l1_conservation_line.png": (
        "RPL22L1 Conservation Score Along Alignment",
        [
            "Per-position conservation score for RPL22L1.",

            "Comparison with the RPL22 conservation plot reveals whether the two paralogs "
            "are under similar or different selective pressures at each position.",
        ]
    ),

    "rpl22_conservation_heatmap.png": (
        "RPL22 Conservation Heatmap",
        [
            "Heatmap where each row is a species and each column is an alignment position. "
            "Blue indicates matching the consensus, pink indicates divergence, light gray "
            "indicates gaps.",

            "Vertical bands of consistent blue reveal globally conserved positions.",
        ]
    ),

    "rpl22l1_conservation_heatmap.png": (
        "RPL22L1 Conservation Heatmap",
        [
            "Conservation heatmap for RPL22L1, using the same color scheme as RPL22.",
        ]
    ),

    "rpl22_identity_matrix.png": (
        "RPL22 Pairwise Sequence Identity Matrix",
        [
            "Percent sequence identity between every pair of species for RPL22.",

            "Closely related species share high identity, while distant comparisons show "
            "lower but still substantial identity, reflecting the essential nature of this "
            "ribosomal protein.",
        ]
    ),

    "rpl22l1_identity_matrix.png": (
        "RPL22L1 Pairwise Sequence Identity Matrix",
        [
            "Pairwise identity matrix for RPL22L1.",

            "Comparing these values with the RPL22 matrix reveals whether the paralog "
            "is under similar or relaxed selective constraint.",
        ]
    ),

    "rpl22_domain_conservation.png": (
        "RPL22 Domain-Specific Conservation",
        [
            "Bar chart comparing mean conservation within the L22e domain versus the "
            "N-terminal and C-terminal extensions.",

            "The L22e domain shows significantly higher conservation than the flanking "
            "regions, consistent with its essential role in ribosome assembly.",
        ]
    ),

    "rpl22l1_domain_conservation.png": (
        "RPL22L1 Domain-Specific Conservation",
        [
            "Domain-specific conservation for RPL22L1.",
        ]
    ),

    "rpl22_conserved_regions_map.png": (
        "RPL22 Conserved Regions Map",
        [
            "Linear map showing conserved regions (>=70%) on the RPL22 protein.",
        ]
    ),

    "rpl22l1_conserved_regions_map.png": (
        "RPL22L1 Conserved Regions Map",
        [
            "Linear map showing conserved regions (>=70%) on the RPL22L1 protein.",
        ]
    ),

    # ── Comparison figures ────────────────────────────────────────────────
    "rpl22_vs_rpl22l1_pairwise.png": (
        "Human RPL22 vs RPL22L1 Pairwise Alignment",
        [
            "ClustalW-style pairwise alignment of human RPL22 and RPL22L1 protein sequences.",

            "This alignment directly compares the two paralogs, highlighting identical "
            "positions (colored backgrounds) and divergent positions. Conservation symbols "
            "indicate the degree of similarity at each position.",

            "Regions of high identity between the paralogs correspond to the core L22e "
            "domain fold required for ribosome incorporation. Divergent regions, particularly "
            "at the N-terminus, may contribute to paralog-specific functions or regulation.",

            "RPL22 directly represses RPL22L1 by binding an internal stem-loop in the "
            "RPL22L1 mRNA. The protein-level differences shown here may affect ribosome "
            "composition and translational specificity.",
        ]
    ),

    "rpl22_vs_rpl22l1_divergence_map.png": (
        "RPL22 vs RPL22L1 Residue-by-Residue Divergence Map",
        [
            "Visual map showing identical (blue), similar (orange), and different (red) "
            "positions between human RPL22 and RPL22L1.",

            "This bar chart provides a quick overview of where the two paralogs have "
            "diverged. Clusters of red bars indicate regions of significant divergence "
            "that may underlie functional differences between the paralogs.",

            "The N-terminal region typically shows more divergence, while the core L22e "
            "domain is more conserved between the paralogs.",
        ]
    ),

    "rpl22_vs_rpl22l1_conservation.png": (
        "RPL22 vs RPL22L1 Cross-Species Conservation Comparison",
        [
            "Overlay of per-position conservation scores for RPL22 (blue) and RPL22L1 (red) "
            "across all species, plotted against human residue positions.",

            "This comparison reveals whether the two paralogs are under similar evolutionary "
            "pressure at each position. Positions where both curves are high indicate residues "
            "essential for the shared ribosomal function.",

            "Positions where RPL22 conservation exceeds RPL22L1 (or vice versa) may indicate "
            "paralog-specific functional constraints or relaxed selection following gene "
            "duplication.",
        ]
    ),
}


SLIDE_GROUPS = [
    ("Overview", [
        "rpl22_combined_conservation_architecture.png",
        "rpl22l1_combined_conservation_architecture.png",
        "rpl22_animal_domain_architecture.png",
        "rpl22l1_animal_domain_architecture.png",
    ]),
    ("RPL22 Conservation Analysis", [
        "rpl22_conservation_line.png",
        "rpl22_conservation_heatmap.png",
        "rpl22_identity_matrix.png",
        "rpl22_domain_conservation.png",
        "rpl22_conserved_regions_map.png",
    ]),
    ("RPL22L1 Conservation Analysis", [
        "rpl22l1_conservation_line.png",
        "rpl22l1_conservation_heatmap.png",
        "rpl22l1_identity_matrix.png",
        "rpl22l1_domain_conservation.png",
        "rpl22l1_conserved_regions_map.png",
    ]),
    ("RPL22 Publication Alignments", [
        "rpl22_pub_full_alignment.png",
        "rpl22_pub_l22e_domain.png",
        "rpl22_pub_domain_architecture.png",
        "rpl22_pub_combined_panels.png",
    ]),
    ("RPL22L1 Publication Alignments", [
        "rpl22l1_pub_full_alignment.png",
        "rpl22l1_pub_l22e_domain.png",
        "rpl22l1_pub_domain_architecture.png",
        "rpl22l1_pub_combined_panels.png",
    ]),
    ("RPL22 vs RPL22L1 Comparison", [
        "rpl22_vs_rpl22l1_pairwise.png",
        "rpl22_vs_rpl22l1_divergence_map.png",
        "rpl22_vs_rpl22l1_conservation.png",
    ]),
]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.25)

SECTION_DESCRIPTIONS = {
    "Overview": "Key summary figures of RPL22 and RPL22L1 conservation across animal species",
    "RPL22 Conservation Analysis": "Quantitative conservation analysis of RPL22 (eL22)",
    "RPL22L1 Conservation Analysis": "Quantitative conservation analysis of RPL22L1 (eL22-like)",
    "RPL22 Publication Alignments": "ClustalW-style alignments and domain architecture for RPL22",
    "RPL22L1 Publication Alignments": "ClustalW-style alignments and domain architecture for RPL22L1",
    "RPL22 vs RPL22L1 Comparison": "Direct comparison of the two ribosomal protein paralogs",
}


def add_section_slide(prs, section_title, section_desc):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(2.0), SLIDE_W - Inches(2), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    if section_desc:
        p2 = tf.add_paragraph()
        p2.text = section_desc
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(16)


def add_figure_slide(prs, img_path, title, explanation_paragraphs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_h = Inches(0.65)
    txTitle = slide.shapes.add_textbox(
        MARGIN, Inches(0.08), SLIDE_W - 2 * MARGIN, title_h)
    tf = txTitle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    p.alignment = PP_ALIGN.LEFT

    with Image.open(img_path) as img:
        img_w, img_h = img.size

    text_panel_w = Inches(4.0)
    img_area_w = SLIDE_W - text_panel_w - 3 * MARGIN
    img_area_h = SLIDE_H - title_h - Inches(0.3)

    scale_w = img_area_w / img_w
    scale_h = img_area_h / img_h
    scale = min(scale_w, scale_h)

    final_w = int(img_w * scale)
    final_h = int(img_h * scale)

    img_left = MARGIN
    img_top = int(title_h + (img_area_h - final_h) / 2)

    slide.shapes.add_picture(img_path, img_left, img_top, final_w, final_h)

    text_left = SLIDE_W - text_panel_w - MARGIN
    text_top = title_h + Inches(0.1)
    text_h = SLIDE_H - title_h - Inches(0.4)

    txBox = slide.shapes.add_textbox(text_left, text_top, text_panel_w, text_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    for idx, para_text in enumerate(explanation_paragraphs):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para_text
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(8) if idx > 0 else Pt(0)
        p.space_after = Pt(4)
        p.alignment = PP_ALIGN.LEFT


def get_explanation(fname):
    if fname in EXPLANATIONS:
        return EXPLANATIONS[fname]
    return (fname.replace(".png", "").replace("_", " ").title(), [])


def main():
    print("=" * 60)
    print("Creating RPL22/RPL22L1 Analysis PowerPoint")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = RGBColor(0x0D, 0x47, 0xA1)

    txBox = title_slide.shapes.add_textbox(
        Inches(1.5), Inches(1.5), SLIDE_W - Inches(3), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "RPL22 and RPL22L1 Sequence Conservation Analysis"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Cross-species alignment and paralog comparison of ribosomal protein eL22"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "Species panel spanning Mammalia to Cnidaria"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(10)

    p4 = tf.add_paragraph()
    p4.text = (
        "RPL22 (eL22): 60S ribosomal protein, EBV EBER-binding, cell cycle regulation\n"
        "RPL22L1 (eL22-like): Paralog that compensates for RPL22 loss in cancer"
    )
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(20)

    p5 = tf.add_paragraph()
    p5.text = (
        "Methods: Sequences from NCBI (RefSeq), aligned with Clustal Omega (EBI), "
        "analyzed with Python (Biopython, Matplotlib)"
    )
    p5.font.size = Pt(11)
    p5.font.color.rgb = RGBColor(0x64, 0xB5, 0xF6)
    p5.alignment = PP_ALIGN.CENTER
    p5.space_before = Pt(10)

    slide_count = 1
    for section_title, figures in SLIDE_GROUPS:
        desc = SECTION_DESCRIPTIONS.get(section_title, "")
        add_section_slide(prs, section_title, desc)
        slide_count += 1
        print(f"\n  Section: {section_title}")

        for fname in figures:
            img_path = os.path.join(OUTPUT_DIR, fname)
            if not os.path.exists(img_path):
                print(f"    SKIP (not found): {fname}")
                continue
            title, paragraphs = get_explanation(fname)
            add_figure_slide(prs, img_path, title, paragraphs)
            slide_count += 1
            print(f"    Added: {fname}")

    out_path = os.path.join(OUTPUT_DIR, "RPL22_RPL22L1_Analysis_Figures.pptx")
    prs.save(out_path)
    print(f"\n{'='*60}")
    print(f"Saved: {out_path}")
    print(f"Total slides: {slide_count}")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()

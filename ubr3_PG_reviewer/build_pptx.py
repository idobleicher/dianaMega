#!/usr/bin/env python
"""Assemble the UBR3 panel deck: one slide per figure panel, large, with a
detailed written explanation beside it.

Layout adapts to panel shape - wide panels get the image across the top with the
explanation in two columns beneath; taller panels get image left, text right.
"""
import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import pptx_text as T
import ubr3_core as U

PANELS_DIR = os.path.join(U.HERE, 'panels')
FIGS_DIR = U.FIGS
OUT = os.path.join(U.HERE, 'UBR3_figures_explained.pptx')

W, H = Inches(13.333), Inches(7.5)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
PAGE = RGBColor(0xF9, 0xF9, 0xF7)
FONT = 'Segoe UI'

WIDE_ASPECT = 1.9          # panels wider than this go top/bottom instead of side-by-side


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAGE
    return s


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, bold=False, color=INK, space_before=0, space_after=4,
         first=False, align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return p


def sections(tf, secs, head_size=11.5, body_size=10.5, first=True):
    """Render (heading, body) pairs into a text frame."""
    for i, (head, body) in enumerate(secs):
        para(tf, head.upper(), head_size - 2.0, bold=True, color=ACCENT,
             space_before=0 if (first and i == 0) else 9, space_after=2,
             first=first and i == 0)
        para(tf, body, body_size, color=INK2, space_after=0)


def fit(img_path, max_w, max_h):
    """Largest (w, h) in EMU preserving aspect."""
    iw, ih = Image.open(img_path).size
    scale = min(max_w / iw, max_h / ih)
    return Emu(int(iw * scale)), Emu(int(ih * scale))


def rule(slide, l, t, w, color=ACCENT, h=Pt(2.5)):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


# ------------------------------------------------------------------ slides
def title_slide(prs):
    s = blank(prs)
    rule(s, Inches(0.9), Inches(2.05), Inches(1.6))
    tf = textbox(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.6))
    para(tf, 'UBR3 N24mer screen', 30, bold=True, color=INK, first=True, space_after=6)
    para(tf, 'Pro/Gly N-termini and the [P/G]-[E/D] motif', 30, bold=True, color=ACCENT,
         space_after=0)
    tf2 = textbox(s, Inches(0.9), Inches(4.1), Inches(10.4), Inches(2.0))
    para(tf2, 'Every panel on its own slide, with a detailed explanation of what is '
              'plotted, how it was computed, what it shows, and what to watch out for.',
         14, color=INK2, first=True, space_after=10)
    para(tf2, f'{len(T.PANELS)} panels in {len(T.FIGURES)} groups  ·  '
              'Source: Supplemental Data 1.xlsx  ·  Companion workbook: '
              'UBR3_PG_substrate_tables.xlsx', 11, color=MUTED)


def definitions_slide(prs):
    s = blank(prs)
    tf = textbox(s, Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.5))
    para(tf, 'Definitions and methods', 22, bold=True, color=INK, first=True)
    rule(s, Inches(0.7), Inches(1.02), Inches(1.3))
    half = (len(T.DEFINITIONS) + 1) // 2
    for col, chunk in enumerate([T.DEFINITIONS[:half], T.DEFINITIONS[half:]]):
        tf = textbox(s, Inches(0.7 + col * 6.25), Inches(1.35), Inches(5.8), Inches(5.7))
        sections(tf, chunk, head_size=11, body_size=9.8)


def figure_slide(prs, key):
    ttl, blurb = T.FIGURES[key]
    s = blank(prs)
    tf = textbox(s, Inches(0.7), Inches(0.42), Inches(12.0), Inches(0.55))
    para(tf, ttl, 20, bold=True, color=INK, first=True)
    rule(s, Inches(0.7), Inches(1.0), Inches(1.3))
    tf = textbox(s, Inches(0.7), Inches(1.3), Inches(12.0), Inches(0.9))
    para(tf, blurb, 12.5, color=INK2, first=True)
    # the composed figure as a reference thumbnail
    name = {
        '2': 'Figure2_motif_necessary_not_sufficient',
        '3': 'Figure3_sequence_logos', '4': 'Figure4_AA_class_analysis',
        '5': 'Figure5_position_residue_heatmaps', '6': 'Figure6_PSI_baseline_stability',
        '9': 'Figure9_stability_x_substrate_classification',
        '10': 'Figure10_downstream_of_motif',
        '11': 'Figure11_PG_substrates_vs_library'}[key]
    path = os.path.join(FIGS_DIR, name + '.png')
    if os.path.exists(path):
        w, h = fit(path, Inches(11.2), Inches(4.75))
        s.shapes.add_picture(path, Emu(int((W - w) / 2)), Inches(2.18), w, h)
    tf = textbox(s, Inches(0.7), Inches(7.05), Inches(12.0), Inches(0.3))
    para(tf, 'The composed panel set. The following slides show each panel full size.',
         9.5, color=MUTED, italic=True, first=True)


def panel_slide(prs, key):
    ttl, secs = T.PANELS[key]
    fig_no, letter = key[:-1], key[-1]     # keys may be two digits, e.g. '10A'
    path = os.path.join(PANELS_DIR, f'Fig{key}.png')
    iw, ih = Image.open(path).size
    wide = (iw / ih) > WIDE_ASPECT

    s = blank(prs)
    tf = textbox(s, Inches(0.55), Inches(0.30), Inches(12.3), Inches(0.75))
    p = tf.paragraphs[0]
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = f'{fig_no}{letter}'
    r.font.size, r.font.bold, r.font.name = Pt(18), True, FONT
    r.font.color.rgb = ORANGE
    r = p.add_run()
    r.text = f'   ·   {ttl}'
    r.font.size, r.font.bold, r.font.name = Pt(18), True, FONT
    r.font.color.rgb = INK
    rule(s, Inches(0.55), Inches(0.98), Inches(12.25), color=RGBColor(0xE1, 0xE0, 0xD9), h=Pt(1.2))

    if wide:
        w, h = fit(path, Inches(12.3), Inches(3.95))
        s.shapes.add_picture(path, Emu(int((W - w) / 2)), Inches(1.18), w, h)
        top = Inches(1.18) + h + Inches(0.16)
        half = (len(secs) + 1) // 2
        for col, chunk in enumerate([secs[:half], secs[half:]]):
            tf = textbox(s, Inches(0.55 + col * 6.35), top, Inches(5.9),
                         H - top - Inches(0.25))
            sections(tf, chunk, head_size=10.5, body_size=9.6)
    else:
        w, h = fit(path, Inches(7.75), Inches(5.95))
        s.shapes.add_picture(path, Inches(0.5), Inches(1.20) + Emu(int((Inches(5.95) - h) / 2)),
                             w, h)
        tf = textbox(s, Inches(8.62), Inches(1.22), Inches(4.25), Inches(5.95))
        sections(tf, secs, head_size=11, body_size=10.2)


def closing_slide(prs):
    s = blank(prs)
    tf = textbox(s, Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.5))
    para(tf, 'Summary of conclusions', 22, bold=True, color=INK, first=True)
    rule(s, Inches(0.7), Inches(1.02), Inches(1.3))
    half = (len(T.CLOSING) + 1) // 2
    for col, chunk in enumerate([T.CLOSING[:half], T.CLOSING[half:]]):
        tf = textbox(s, Inches(0.7 + col * 6.25), Inches(1.35), Inches(5.8), Inches(5.7))
        sections(tf, chunk, head_size=11.5, body_size=10.4)


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    title_slide(prs)
    definitions_slide(prs)
    n = 2
    for fig in sorted(T.FIGURES, key=int):
        figure_slide(prs, fig)
        n += 1
        for key in sorted((k for k in T.PANELS if k[:-1] == fig), key=lambda k: k[-1]):
            panel_slide(prs, key)
            n += 1
    closing_slide(prs)
    n += 1

    prs.save(OUT)
    print(f'wrote {OUT}  ({n} slides)')


if __name__ == '__main__':
    main()

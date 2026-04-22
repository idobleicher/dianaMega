"""
Build a publication-ready PowerPoint deck that organizes all the figures
produced in this analysis, with per-slide explanations, the Methods
section, raw data tables, and references.
"""
import os, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image

# ══════════════════════════════════════════════════════════════════════════
#  CONFIG
# ══════════════════════════════════════════════════════════════════════════
OUT_DIR  = r'c:\Users\User\Desktop\תינוקת\dianaMega\ubr3enrichmentlogo'
OUT_FILE = os.path.join(OUT_DIR, 'UBR3_N_terminal_screen_paper_deck.pptx')

FIG_DIR  = r'c:\Users\User\Desktop\תינוקת\dianaMega\ubr3enrichmentlogo\figures'
DIP_DIR  = r'c:\Users\User\Desktop\תינוקת\dianaMega\ubr3enrichmentlogo\dipeptide_analysis'

COLOR_TITLE       = RGBColor(0x78, 0x28, 0x1F)
COLOR_ACCENT      = RGBColor(0xC0, 0x39, 0x2B)
COLOR_SUBTLE      = RGBColor(0x55, 0x55, 0x55)
COLOR_DARK        = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_BG_METHODS  = RGBColor(0xFD, 0xF6, 0xE3)
COLOR_BG_DATA     = RGBColor(0xEA, 0xF2, 0xF8)
COLOR_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════════
#  HELPERS
# ══════════════════════════════════════════════════════════════════════════
def add_rect(slide, left, top, w, h, fill_rgb, line_rgb=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp

def add_textbox(slide, left, top, w, h, text, *,
                size=14, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        if color is not None:
            run.font.color.rgb = color
    return tb

def add_bullets(slide, left, top, w, h, items, *,
                size=14, color=None, bullet_char='•', line_space=1.15):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_space
        run = p.add_run()
        run.text = f'{bullet_char}  {item}'
        run.font.name = 'Calibri'
        run.font.size = Pt(size)
        if color is not None:
            run.font.color.rgb = color
    return tb

def add_image_scaled(slide, png_path, left, top, max_w, max_h):
    im = Image.open(png_path)
    w_px, h_px = im.size
    dpi = im.info.get('dpi', (300, 300))[0]
    src_w_in = w_px / dpi
    src_h_in = h_px / dpi
    ratio = min(max_w / Inches(src_w_in), max_h / Inches(src_h_in))
    out_w = int(Inches(src_w_in) * ratio)
    out_h = int(Inches(src_h_in) * ratio)
    offset_l = left + (max_w - out_w) // 2
    offset_t = top + (max_h - out_h) // 2
    slide.shapes.add_picture(png_path, offset_l, offset_t, out_w, out_h)

def add_header_bar(slide, text, subtitle=None, color=COLOR_TITLE):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.75), color)
    add_textbox(slide, Inches(0.4), Inches(0.1), SLIDE_W - Inches(0.8),
                Inches(0.55), text,
                size=22, bold=True, color=COLOR_WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(0.77),
                    SLIDE_W - Inches(0.8), Inches(0.35),
                    subtitle, size=12, italic=True, color=COLOR_SUBTLE)

def add_footer(slide, page, total):
    add_textbox(slide, Inches(0.3), SLIDE_H - Inches(0.35),
                SLIDE_W - Inches(0.6), Inches(0.25),
                f'UBR3 N-terminal screen  ·  {page}/{total}',
                size=9, italic=True, color=COLOR_SUBTLE,
                align=PP_ALIGN.RIGHT)

def add_table(slide, left, top, w, h, data, *,
              header_fill=COLOR_DARK, header_color=COLOR_WHITE,
              font_size=10, first_col_bold=True):
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, w, h)
    table = tbl_shape.table
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.margin_left   = Inches(0.05)
            cell.margin_right  = Inches(0.05)
            cell.margin_top    = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.name = 'Calibri'
            run.font.size = Pt(font_size)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                run.font.bold = True
                run.font.color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE
                if first_col_bold and c == 0:
                    run.font.bold = True
    return tbl_shape

# ══════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════
slides = []

# ── Slide 1: TITLE ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
slides.append(s)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, COLOR_TITLE)
add_rect(s, 0, Inches(5.8), SLIDE_W, Inches(0.08), COLOR_ACCENT)
add_textbox(s, Inches(0.8), Inches(2.2), SLIDE_W - Inches(1.6), Inches(1.3),
            'UBR3 N-terminal Degron Screen',
            size=44, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.8), Inches(3.5), SLIDE_W - Inches(1.6), Inches(0.7),
            'Statistical Reanalysis, Figures and Methods',
            size=26, italic=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.8), Inches(4.5), SLIDE_W - Inches(1.6), Inches(0.5),
            'Library: 16,514 peptides    ·    Best hits: 54 peptides',
            size=16, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.8), Inches(6.0), SLIDE_W - Inches(1.6), Inches(0.5),
            'Figures · Raw data · Methods · References',
            size=14, italic=True, color=RGBColor(0xEC, 0xF0, 0xF1),
            align=PP_ALIGN.CENTER)

# ── Slide 2: CONTENTS ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT); slides.append(s)
add_header_bar(s, 'Contents',
               'Each figure is on its own slide with a paper-ready caption. '
               'Methods and raw data are at the end.')
contents = [
    '1.   Amino-acid enrichment heatmap — Fisher exact test   (Fig 1B)',
    '2.   Top dipeptide enrichments at positions 2–3   (Fig 3)',
    '3.   Fold-change sequence logo (positions 2 and 3)   (Fig 4)',
    '4.   Combined dipeptide groups PD+PE+PT and GE+GD — Mann-Whitney U',
    '5.   Combined dipeptide groups — Brunner-Munzel + Cohen\'s d',
    '6.   Volcano plot — moderated t-test (limma-style EB) + BH FDR   (Fig 20c)',
    '7.   Methods section   (paper-ready)',
    '8.   Raw data — Top Hits (PD+PE+PT and GE+GD)',
    '9.   Software, packages, references',
]
add_bullets(s, Inches(0.6), Inches(1.2), SLIDE_W - Inches(1.2), Inches(5.5),
            contents, size=18, color=COLOR_DARK, bullet_char='▸',
            line_space=1.4)
add_footer(s, 2, 0)

# ══════════════════════════════════════════════════════════════════════════
#  FIGURE-SLIDE helper
# ══════════════════════════════════════════════════════════════════════════
def figure_slide(title, subtitle, png_path, description, method_summary,
                 svg_note=True):
    s = prs.slides.add_slide(BLANK_LAYOUT); slides.append(s)
    add_header_bar(s, title, subtitle)
    # Left: figure
    fig_left = Inches(0.4)
    fig_top  = Inches(1.0)
    fig_w    = Inches(7.6)
    fig_h    = Inches(5.8)
    add_rect(s, fig_left, fig_top, fig_w, fig_h, COLOR_WHITE,
             line_rgb=RGBColor(0xD5, 0xD8, 0xDC))
    add_image_scaled(s, png_path,
                     fig_left + Inches(0.08), fig_top + Inches(0.08),
                     fig_w - Inches(0.16), fig_h - Inches(0.16))
    # Right: description + method
    right_left = Inches(8.2)
    right_w    = Inches(4.8)
    add_textbox(s, right_left, Inches(1.05), right_w, Inches(0.35),
                'Description', size=14, bold=True, color=COLOR_ACCENT)
    add_bullets(s, right_left, Inches(1.42), right_w, Inches(2.8),
                description, size=11, color=COLOR_DARK, bullet_char='•',
                line_space=1.25)
    add_textbox(s, right_left, Inches(4.35), right_w, Inches(0.35),
                'Statistical method', size=14, bold=True, color=COLOR_ACCENT)
    add_bullets(s, right_left, Inches(4.72), right_w, Inches(2.0),
                method_summary, size=10.5, color=COLOR_DARK, bullet_char='•',
                line_space=1.2)
    if svg_note:
        add_textbox(s, right_left, Inches(6.7), right_w, Inches(0.3),
                    'Editable vector (SVG) available in project folder',
                    size=9, italic=True, color=COLOR_SUBTLE)

# ── Slide 3: Fig 1B — Enrichment heatmap ──────────────────────────────────
figure_slide(
    'Fig 1B.  Amino-acid enrichment heatmap (Positions 2 and 3)',
    'Best hits vs. full library  |  Fisher exact test',
    os.path.join(FIG_DIR, 'fig1b_enrichment_heatmap.png'),
    description=[
        'Heatmap of per-residue enrichment at the two positions immediately '
        'after the N-terminal Met.',
        'Enrichment = frequency in best hits / frequency in library. '
        'Rows sorted by enrichment at position 2.',
        'P, G and T dominate position 2; D, E and N dominate position 3 — '
        'consistent with the UBR family Nt-degron recognition rules.',
        'Stars indicate significance from Fisher exact tests (two-sided) on '
        '2×2 contingency tables (hit vs. non-hit × carries AA vs. does not).',
    ],
    method_summary=[
        'Test: Fisher exact (two-sided)   scipy.stats.fisher_exact',
        '2×2 table per amino acid × position',
        'n_hits = 54,   n_library = 16,514',
        '*  p<0.05    **  p<0.01    ***  p<0.001',
    ],
)

# ── Slide 4: Fig 3 — Top dipeptide bars ───────────────────────────────────
figure_slide(
    'Fig 3.  Top-10 dipeptide enrichments at P2-P3',
    'Best hits vs. full library  |  frequency fold-change',
    os.path.join(FIG_DIR, 'fig3_dipeptide_enrichment_bars.png'),
    description=[
        'Ranked enrichment of the 10 most frequent dipeptides at positions '
        '2–3 in the best-hits subset (out of the top-20 observed).',
        'Dipeptide enrichment = f(dipeptide|hits) / f(dipeptide|library).',
        'PD, PT, GD and GE exceed 30× enrichment — these are the dominant '
        'P2-P3 motifs in UBR3 substrates.',
        'Dashed line at 1 marks the background (library) frequency.',
    ],
    method_summary=[
        'Frequency-based enrichment, no statistical test (descriptive).',
        'Chart truncated at dipeptide EL (rank 10).',
        'See Fig 1C in the project for the complete 20×20 heatmap.',
    ],
)

# ── Slide 5: Fig 4 — Sequence logo ────────────────────────────────────────
figure_slide(
    'Fig 4.  Fold-change sequence logo',
    'Best hits vs. library at positions 2 and 3',
    os.path.join(FIG_DIR, 'fig4_logo_pos2_pos3_besthits.png'),
    description=[
        'Stacked logo where each letter\'s height equals its fold-change '
        '(enrichment) at that position; letters with FC ≤ 1 are hidden.',
        'Colour = biochemical category: Acidic (dark red), Basic, '
        'Nonpolar (orange) and Polar (yellow).',
        'Position 2 is dominated by P and G (large Nonpolar stack); '
        'position 3 is dominated by D and E (Acidic stack).',
        'A simultaneous P/G at P2 and D/E at P3 defines the core UBR3 motif.',
    ],
    method_summary=[
        'Built with logomaker (Python); stack_order="big_on_top".',
        'Matrix = per-residue enrichment (hits/library); threshold at 1.',
        'No statistical test — descriptive visualization of enrichment.',
    ],
)

# ── Slide 6: Combined boxplot — Mann-Whitney ──────────────────────────────
figure_slide(
    'Combined dipeptide groups — Mann-Whitney U test',
    'PD+PE+PT and GE+GD — PSI in control (AAVS) cells',
    os.path.join(DIP_DIR, 'Fig_Combined_only_boxplot.png'),
    description=[
        'PSI (Protein Stability Index in the AAVS control cell line) compared '
        'between Top Hits and the full library for two pooled motif groups.',
        'PD+PE+PT:  n_hits = 8,  n_all = 104  →  p = 0.016  (*)',
        'GE+GD:     n_hits = 12, n_all = 109  →  p = 0.051  (ns, borderline)',
        'One-tailed test (H₁: Top Hits < All Screen) reflecting the '
        'pre-specified hypothesis that destabilizing motifs have lower PSI.',
    ],
    method_summary=[
        'Test: one-tailed Mann-Whitney U',
        '   scipy.stats.mannwhitneyu(alternative="less")',
        'Rank-based; assumes similar dispersions in both groups.',
        'With unequal SDs + unequal N, MW is conservative here → see next slide.',
    ],
)

# ── Slide 7: Combined boxplot — Brunner-Munzel + Cohen's d ────────────────
figure_slide(
    'Combined dipeptide groups — Brunner-Munzel + Cohen\'s d',
    'Rank-based test robust to unequal variances and unequal n',
    os.path.join(DIP_DIR, 'Fig_Combined_only_cohen_d.png'),
    description=[
        'Same data as the previous slide, analysed with a test appropriate '
        'for the data structure (unequal variances, unequal sample sizes, '
        'non-normal library distribution).',
        'PD+PE+PT:  d = 1.05 (large effect),   Brunner-Munzel p = 0.0017  **',
        'GE+GD:     d = 0.49 (small effect),    Brunner-Munzel p = 0.0086  **',
        'Both dipeptide groups are significantly destabilized relative to the '
        'background screen.',
    ],
    method_summary=[
        'Test: one-tailed Brunner-Munzel',
        '   scipy.stats.brunnermunzel(alternative="less")',
        'Cohen\'s d = (mean_all − mean_hits) / pooled SD.',
        'Rank-based, Welch-style variance adjustment — rank analogue of '
        'Welch\'s t-test.',
        'Shapiro-Wilk on All Screen: non-normal (p < 0.001) → parametric '
        'tests are not ideal.',
    ],
)

# ── Slide 8: Fig 20c — Moderated-t volcano ────────────────────────────────
figure_slide(
    'Fig 20C.  Volcano plot — moderated t-test + BH FDR',
    'limma-style empirical Bayes variance shrinkage (Smyth 2004)',
    os.path.join(FIG_DIR, 'fig20c_volcano_moderated_t.png'),
    description=[
        'ΔPSI = UBR3-cell PSI − AAVS-control PSI, per peptide (N=2 each).',
        'Y-axis = –log10(BH-adjusted p) from a moderated t-test — variances '
        'are shrunk toward a dataset-wide prior, raising effective d.f. from '
        '~1–2 (Welch) to ~11.6.',
        'Empirical Bayes prior:  s₀² = 0.013,  d₀ = 9.6 (fit from all 16,514 variances).',
        'After |ΔPSI| ≥ 0.5 and FDR < 0.05:  43 stabilized, 16 destabilized, '
        'plus the 54 pre-labelled best hits — all visibly enriched in the '
        'upper-right quadrant.',
        'Dots coloured by Position-2 residue; hits labelled Gene (P2P3).',
    ],
    method_summary=[
        'Moderated t-test (Smyth 2004) implemented from scratch.',
        'Per-gene pooled variance s²_g  (d.f.  d_g = 2).',
        'Prior (s₀², d₀) estimated by fitting a scaled F-distribution to s²_g '
        '(trigamma-based method of moments).',
        's̃²_g = (d₀·s₀² + d_g·s²_g) / (d₀ + d_g)',
        't̃_g = ΔPSI / √[s̃²_g · (1/n₁ + 1/n₂)]',
        'p̃_g from t-distribution with (d_g + d₀) d.f., two-sided.',
        'Benjamini-Hochberg FDR:   statsmodels.multitest, method="fdr_bh".',
    ],
)

# ══════════════════════════════════════════════════════════════════════════
#  METHODS SLIDES
# ══════════════════════════════════════════════════════════════════════════
# ── Slide 9: Methods part 1 ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT); slides.append(s)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, COLOR_BG_METHODS)
add_header_bar(s, 'Methods — Data and statistics (paper-ready, part 1)',
               'Verbatim text suitable for the Methods section of the paper.',
               color=COLOR_DARK)

methods_text_1 = (
"Dataset.  A library of 16,514 human N-terminal 24-mer peptides was "
"screened in duplicate in HEK293-derived AAVS (control) and UBR3-expressing "
"cells using a GFP:RFP dual-reporter stability assay. For each peptide, the "
"Protein Stability Index (PSI) was taken as the mean log-ratio over the two "
"replicates: PSI_ctrl = (PSI-293a + PSI-293b) / 2 and "
"PSI_UBR3 = (PSI-UBR3a + PSI-UBR3b) / 2.  "
"The per-peptide differential was ΔPSI = PSI_UBR3 − PSI_ctrl, where a "
"positive value indicates UBR3-dependent stabilization. A curated set of 54 "
"top hits (sheet sub_high, column HITubr3 = TRUE) was used wherever the "
"analysis refers to \"best hits\"."

"\n\nEnrichment analyses (Figs 1B, 3, 4).  For each amino acid at each "
"position and for each dipeptide at positions 2–3, enrichment was defined "
"as f_hits / f_library where f is the frequency in the respective subset. "
"Per-residue significance (Fig 1B) was assessed with a two-sided Fisher "
"exact test on the 2×2 contingency table [hit vs non-hit] × [carries AA at "
"position p vs does not] (scipy.stats.fisher_exact). Sequence logos (Fig 4) "
"were built in logomaker 0.8 using the fold-change matrix with stack_order "
"= \"big_on_top\"; letters with fold-change ≤ 1 were hidden."
)
add_textbox(s, Inches(0.5), Inches(1.1), SLIDE_W - Inches(1.0), Inches(6.0),
            methods_text_1, size=12, color=COLOR_DARK,
            align=PP_ALIGN.JUSTIFY, anchor=MSO_ANCHOR.TOP)
add_footer(s, 9, 0)

# ── Slide 10: Methods part 2 ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT); slides.append(s)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, COLOR_BG_METHODS)
add_header_bar(s, 'Methods — Statistics (part 2)',
               'Group comparisons and moderated t-test.',
               color=COLOR_DARK)

methods_text_2 = (
"Comparison of combined dipeptide groups.  Peptides were pooled into two "
"groups: PD+PE+PT and GE+GD. For each group, the PSI_ctrl distribution of "
"Top Hits was compared with the PSI_ctrl distribution of all library "
"peptides carrying the same motif. Because the two groups in each "
"comparison differed markedly in sample size (n_hits ≈ 8–12 vs n_library "
"≈ 100) and dispersion (SD_library / SD_hits ≈ 1.7), and because the "
"library distribution was non-normal (Shapiro-Wilk p < 0.001), the "
"reported p-values are from the Brunner-Munzel test "
"(scipy.stats.brunnermunzel, alternative=\"less\"), which is the rank "
"analogue of Welch's t-test and does not assume equal variances or equal "
"shape. Effect sizes are Cohen's d with a pooled standard deviation "
"(ddof = 1). The one-tailed Mann-Whitney U test "
"(scipy.stats.mannwhitneyu, alternative=\"less\") is reported as a "
"supplementary reference in the original boxplot figure."

"\n\nDifferential stability (volcano plot, Fig 20C).  To improve power "
"given the N=2 replicate design, ΔPSI was analysed with a moderated "
"t-test using the limma-style empirical Bayes variance shrinkage of Smyth "
"(2004). Per-peptide pooled variances s²_g were computed under the "
"equal-variance assumption (d.f. d_g = 2). A scaled F-distribution was "
"fitted to the 16,514 variances by the method of moments on log(s²_g), "
"using the trigamma-inverse estimator, yielding prior hyperparameters "
"s₀² = 0.0133 and d₀ = 9.61. Shrunk variances were computed as "
"s̃²_g = (d₀·s₀² + d_g·s²_g) / (d₀ + d_g). The moderated t-statistic was "
"t̃_g = ΔPSI / √[s̃²_g · (1/n₁ + 1/n₂)] and p-values were obtained from "
"a t-distribution with d_g + d₀ ≈ 11.6 d.f. Multiple-testing correction "
"was performed with the Benjamini-Hochberg procedure "
"(statsmodels.stats.multitest.multipletests, method=\"fdr_bh\"). "
"A peptide was considered significantly regulated when |ΔPSI| ≥ 0.5 "
"and FDR < 0.05. As an orthogonal, non-parametric cross-check, a robust "
"z-score of ΔPSI was computed as (ΔPSI − median) / (1.4826·MAD)."
)
add_textbox(s, Inches(0.5), Inches(1.1), SLIDE_W - Inches(1.0), Inches(6.0),
            methods_text_2, size=12, color=COLOR_DARK,
            align=PP_ALIGN.JUSTIFY, anchor=MSO_ANCHOR.TOP)
add_footer(s, 10, 0)

# ══════════════════════════════════════════════════════════════════════════
#  RAW DATA SLIDES
# ══════════════════════════════════════════════════════════════════════════
# ── Slide 11: Raw data PD+PE+PT ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT); slides.append(s)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, COLOR_BG_DATA)
add_header_bar(s, 'Raw data — Top Hits with PD, PE or PT at P2-P3',
               'n_hits = 8     |     n_all_library = 104     |     mean PSI = 1.835 vs 2.490',
               color=COLOR_DARK)

pd_data = [
    ['Gene',     'Dipeptide', 'PSI-293a', 'PSI-293b', 'PSI_AAVS'],
    ['UNC13B',   'PD',        '1.4356',   '1.5907',   '1.5132'],
    ['PGBD5',    'PD',        '2.3598',   '2.4766',   '2.4182'],
    ['UVRAG',    'PD',        '1.5629',   '1.3280',   '1.4454'],
    ['ERMAP',    'PE',        '1.7786',   '1.8865',   '1.8326'],
    ['FARSB',    'PT',        '1.7000',   '1.7056',   '1.7028'],
    ['B3GNTL1',  'PT',        '1.3689',   '1.4457',   '1.4073'],
    ['ZNF496',   'PT',        '2.0584',   '1.9931',   '2.0258'],
    ['INTS14',   'PT',        '2.4953',   '2.1736',   '2.3345'],
]
add_table(s, Inches(0.6), Inches(1.1), Inches(7.5), Inches(5.5), pd_data,
          font_size=13)

add_textbox(s, Inches(8.3), Inches(1.1), Inches(4.8), Inches(0.4),
            'Summary statistics', size=14, bold=True, color=COLOR_ACCENT)
summary_pd = [
    'n hits       = 8',
    'mean PSI   = 1.8350',
    'SD PSI      = 0.3931',
    'median      = 1.7677',
    'Q1, Q3       = 1.496, 2.103',
    'vs all-library PSI  (n = 104):',
    '   mean       = 2.4900',
    '   SD          = 0.7907',
    'mean diff      = −0.655',
    "Cohen's d     =  1.05  (large)",
    'Brunner-Munzel  p = 0.0017',
]
add_bullets(s, Inches(8.3), Inches(1.55), Inches(4.8), Inches(5.0),
            summary_pd, size=12, color=COLOR_DARK, bullet_char='·')
add_footer(s, 11, 0)

# ── Slide 12: Raw data GE+GD ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT); slides.append(s)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, COLOR_BG_DATA)
add_header_bar(s, 'Raw data — Top Hits with GE or GD at P2-P3',
               'n_hits = 12     |     n_all_library = 109     |     mean PSI = 2.562 vs 2.836',
               color=COLOR_DARK)

ge_data = [
    ['Gene',     'Dipeptide', 'PSI-293a', 'PSI-293b', 'PSI_AAVS'],
    ['FARP1',    'GE',        '2.4875',   '2.6695',   '2.5785'],
    ['C17orf51', 'GE',        '2.0616',   '2.0534',   '2.0575'],
    ['FBXO24',   'GE',        '1.9973',   '2.1763',   '2.0868'],
    ['FARP2',    'GE',        '2.4798',   '2.4433',   '2.4616'],
    ['KIAA1211', 'GE',        '2.2567',   '2.2608',   '2.2587'],
    ['LPIN1',    'GE',        '2.8316',   '2.7794',   '2.8055'],
    ['MYLK',     'GD',        '2.7720',   '2.5882',   '2.6801'],
    ['LTBP4',    'GD',        '2.1862',   '2.1254',   '2.1558'],
    ['SLC25A4',  'GD',        '2.3997',   '2.4494',   '2.4246'],
    ['ATF7',     'GD',        '2.9843',   '2.8472',   '2.9158'],
    ['SLC7A9',   'GD',        '3.1908',   '3.3571',   '3.2740'],
    ['THEG',     'GD',        '3.2440',   '2.8405',   '3.0423'],
]
add_table(s, Inches(0.6), Inches(1.1), Inches(7.5), Inches(5.8), ge_data,
          font_size=11)

add_textbox(s, Inches(8.3), Inches(1.1), Inches(4.8), Inches(0.4),
            'Summary statistics', size=14, bold=True, color=COLOR_ACCENT)
summary_ge = [
    'n hits       = 12',
    'mean PSI   = 2.5618',
    'SD PSI      = 0.3934',
    'median      = 2.5200',
    'Q1, Q3       = 2.233, 2.833',
    'vs all-library PSI  (n = 109):',
    '   mean       = 2.8364',
    '   SD          = 0.6937',
    'mean diff      = −0.275',
    "Cohen's d     =  0.49  (small)",
    'Brunner-Munzel  p = 0.0086',
]
add_bullets(s, Inches(8.3), Inches(1.55), Inches(4.8), Inches(5.0),
            summary_ge, size=12, color=COLOR_DARK, bullet_char='·')
add_footer(s, 12, 0)

# ══════════════════════════════════════════════════════════════════════════
#  SOFTWARE / REFERENCES
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT); slides.append(s)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, COLOR_BG_METHODS)
add_header_bar(s, 'Software and references', '', color=COLOR_DARK)

add_textbox(s, Inches(0.5), Inches(1.0), Inches(6.0), Inches(0.4),
            'Python environment', size=16, bold=True, color=COLOR_ACCENT)
software_list = [
    'Python 3.13',
    'numpy 2.4, pandas 2.3, scipy 1.17',
    'matplotlib 3.10',
    'logomaker 0.8         (Tareen & Kinney 2019)',
    'statsmodels 0.14      (BH FDR correction)',
    'adjustText 1.3         (non-overlapping labels)',
    'python-docx 1.2, python-pptx 1.0, PyMuPDF 1.27',
    'openpyxl / xlsxwriter (Excel I/O)',
]
add_bullets(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(3.0),
            software_list, size=12, color=COLOR_DARK, bullet_char='·')

add_textbox(s, Inches(0.5), Inches(4.5), Inches(6.0), Inches(0.4),
            'Key statistical references', size=16, bold=True, color=COLOR_ACCENT)
refs_list = [
    'Smyth GK. Linear models and empirical Bayes methods for assessing '
    'differential expression in microarray experiments. Stat Appl Genet Mol '
    'Biol. 2004;3:Article 3.',
    'Brunner E, Munzel U. The nonparametric Behrens-Fisher problem: '
    'Asymptotic theory and a small-sample approximation. Biom J. '
    '2000;42:17-25.',
    'Benjamini Y, Hochberg Y. Controlling the false discovery rate: A '
    'practical and powerful approach to multiple testing. J R Stat Soc B. '
    '1995;57:289-300.',
    'Tareen A, Kinney JB. Logomaker: beautiful sequence logos in Python. '
    'Bioinformatics. 2020;36:2272-2274.',
]
add_bullets(s, Inches(0.5), Inches(4.9), Inches(6.0), Inches(2.3),
            refs_list, size=10.5, color=COLOR_DARK, bullet_char='·',
            line_space=1.25)

add_textbox(s, Inches(7.0), Inches(1.0), Inches(6.0), Inches(0.4),
            'Files produced by this analysis', size=16, bold=True,
            color=COLOR_ACCENT)
files_list = [
    'figures/fig1b_enrichment_heatmap.{png|pdf|svg}',
    'figures/fig3_dipeptide_enrichment_bars.{png|pdf|svg}',
    'figures/fig4_logo_pos2_pos3_besthits.{png|pdf|svg}',
    'dipeptide_analysis/Fig_Combined_only_boxplot.{png|pdf|svg}',
    'dipeptide_analysis/Fig_Combined_only_cohen_d.{png|pdf|svg}',
    'figures/fig20c_volcano_moderated_t.{png|pdf|svg}',
    'figures/volcano_moderated_t_results.xlsx',
    'dipeptide_analysis/RAW_DATA_for_verification.xlsx',
    'dipeptide_analysis/Statistical_reanalysis_report.docx',
]
add_bullets(s, Inches(7.0), Inches(1.4), Inches(6.0), Inches(4.0),
            files_list, size=10.5, color=COLOR_DARK, bullet_char='·',
            line_space=1.2)

add_textbox(s, Inches(7.0), Inches(5.3), Inches(6.0), Inches(0.4),
            'Source scripts', size=16, bold=True, color=COLOR_ACCENT)
scripts_list = [
    'ubr3enrichmentlogo/peptide_analysis.py',
    'ubr3paper/dipeptide_boxplot_analysis.py',
    'ubr3paper/volcano_moderated_t.py',
    'ubr3paper/build_paper_pptx.py    (this deck)',
]
add_bullets(s, Inches(7.0), Inches(5.7), Inches(6.0), Inches(1.5),
            scripts_list, size=10.5, color=COLOR_DARK, bullet_char='·',
            line_space=1.2)
add_footer(s, 13, 0)

# ══════════════════════════════════════════════════════════════════════════
#  BACK COVER
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT); slides.append(s)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, COLOR_TITLE)
add_textbox(s, Inches(0.8), Inches(3.0), SLIDE_W - Inches(1.6), Inches(1.0),
            'End of deck',
            size=40, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.8), Inches(4.0), SLIDE_W - Inches(1.6), Inches(0.6),
            'All figures, raw data, and source scripts are kept in the '
            'project repository alongside this deck.',
            size=14, italic=True, color=RGBColor(0xEC, 0xF0, 0xF1),
            align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
#  SAVE + fix footer page counts
# ══════════════════════════════════════════════════════════════════════════
# (Re-edit footer strings now that we know the total.)
total = len(slides)
# Rebuild page numbers on slides that have a footer
# Scan and update the footer text boxes' last line.
for i, slide in enumerate(slides, start=1):
    for shp in slide.shapes:
        if shp.has_text_frame:
            txt = shp.text_frame.text
            if 'UBR3 N-terminal screen' in txt and '/' in txt and '·' in txt:
                for para in shp.text_frame.paragraphs:
                    for run in para.runs:
                        if 'UBR3 N-terminal screen' in run.text:
                            run.text = f'UBR3 N-terminal screen  ·  {i}/{total}'

prs.save(OUT_FILE)
print(f"Saved deck: {OUT_FILE}")
print(f"   {total} slides")

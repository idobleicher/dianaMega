"""
Builds the slide deck for the UBR3 P-[D/E] enrichment analysis.

Output: 2A_UBR3_enrichment.pptx
"""
import json
import pathlib

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = pathlib.Path(__file__).parent

INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x8A, 0x88, 0x80)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
GREEN = RGBColor(0x1B, 0xAF, 0x7A)
AMBER = RGBColor(0xB0, 0x6D, 0x00)
LIGHT = RGBColor(0xF2, 0xF2, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

res = pd.read_csv(HERE / "enrichment_results.csv")
summ = json.loads((HERE / "enrichment_summary.json").read_text())
PRIMARY = summ["primary_foreground"]
red = summ["redundancy"]
prim = res[res["foreground"] == PRIMARY].set_index("background")
raw = res[res["foreground"] == "All 2A instances (raw)"].set_index("background")
B_ALL = "All amino acids, all positions"
B_P2 = "Position 2 of viral proteins"
B_PP = "Residue after any proline (P+1)"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
W = prs.slide_width


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size=16, bold=False, color=INK, space_after=8, first=False,
         align=PP_ALIGN.LEFT, italic=False, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return p


def header(s, title, sub=None):
    tf = box(s, 0.6, 0.35, 12.2, 1.0)
    para(tf, title, size=30, bold=True, color=INK, first=True, space_after=2)
    if sub:
        para(tf, sub, size=14, color=INK2, space_after=0)
    bar = s.shapes.add_shape(1, Inches(0.6), Inches(1.32), Inches(1.5), Pt(3.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False


def footer(s, text):
    tf = box(s, 0.6, 6.95, 12.2, 0.4)
    para(tf, text, size=10, color=MUTED, first=True, space_after=0)


def panel(s, x, y, w, h, fill=LIGHT):
    sh = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def stat(s, x, y, w, value, label, color=BLUE, vsize=34):
    panel(s, x, y, w, 1.5)
    tf = box(s, x + 0.15, y + 0.14, w - 0.3, 1.25)
    para(tf, value, size=vsize, bold=True, color=color, first=True, space_after=2,
         align=PP_ALIGN.CENTER)
    para(tf, label, size=11.5, color=INK2, space_after=0, align=PP_ALIGN.CENTER)


# ============================================================ 1 title
s = slide()
bg = s.shapes.add_shape(1, 0, 0, W, Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE
bg.line.fill.background()
bg.shadow.inherit = False
accent = s.shapes.add_shape(1, 0, 0, Inches(0.22), Inches(7.5))
accent.fill.solid()
accent.fill.fore_color.rgb = BLUE
accent.line.fill.background()
accent.shadow.inherit = False

tf = box(s, 0.9, 2.0, 11.6, 3.4)
para(tf, "Is the UBR3 P–[D/E] motif enriched", size=40, bold=True, color=INK,
     first=True, space_after=2)
para(tf, "after the 2A ribosomal skip?", size=40, bold=True, color=BLUE, space_after=18)
para(tf, "Amino-acid enrichment at position 2 of 2A downstream products,",
     size=17, color=INK2, space_after=2)
para(tf, "tested against three backgrounds built from 17,517 viral proteins",
     size=17, color=INK2, space_after=20)
para(tf, f"Foreground: Rao et al. 2025 combined 2A set · "
         f"{red['instances']:,} instances → {red['unique_downstream_contexts']:,} "
         f"distinct downstream contexts", size=12, color=MUTED, space_after=0)

# ============================================================ 2 the biology
s = slide()
header(s, "The question", "What the 2A skip produces, and what UBR3 reads")
tf = box(s, 0.6, 1.7, 6.2, 4.9)
para(tf, "The 2A ribosomal skip", size=19, bold=True, color=BLUE, first=True,
     space_after=8)
para(tf, "2A peptides end in a conserved …D(V/I)ExNPG↓P motif. The ribosome skips "
         "forming the peptide bond between the terminal glycine and the following "
         "proline, releasing the upstream protein and continuing translation.",
     size=15, color=INK, space_after=12)
para(tf, "The consequence: every downstream product begins with that proline. "
         "Position 1 is P, always — it is built into the mechanism.",
     size=15, color=INK, space_after=16)
para(tf, "What UBR3 recognises", size=19, bold=True, color=BLUE, space_after=8)
para(tf, "UBR3 is an N-recognin E3 ligase. The motif of interest is a proline at "
         "position 1 followed by an acidic residue — aspartate (D) or glutamate (E) — "
         "at position 2.", size=15, color=INK, space_after=12)
para(tf, "Position 2 is therefore the only variable position, and it is contributed by "
         "the protein downstream of the skip site — not by the 2A peptide itself. "
         "That fact drives the whole analysis.", size=15, color=INK, space_after=0)

panel(s, 7.1, 1.75, 5.7, 2.5)
tf = box(s, 7.35, 1.95, 5.2, 2.2)
para(tf, "…  D V E x N P G   ↓   P  ·  X  …", size=22, bold=True, color=INK,
     first=True, space_after=10, align=PP_ALIGN.CENTER)
para(tf, "skip happens here  →  P is position 1", size=13, color=INK2, space_after=6,
     align=PP_ALIGN.CENTER)
para(tf, "X is position 2 — D or E = UBR3 motif", size=13, bold=True, color=AMBER,
     space_after=0, align=PP_ALIGN.CENTER)

tf = box(s, 7.1, 4.5, 5.7, 2.2)
para(tf, "The question in one line", size=17, bold=True, color=INK, first=True,
     space_after=8)
para(tf, "Does D/E turn up at position 2 more often than chance would put it there — "
         "i.e. is the 2A skip site biased toward or against creating UBR3 substrates?",
     size=15, color=INK, space_after=0)
footer(s, "2A definition and dataset: Rao et al. (2025) Cell Reports, "
          "doi:10.1016/j.celrep.2025.115822")

# ============================================================ 3 data & backgrounds
s = slide()
header(s, "Data and the three backgrounds",
       "A percentage means nothing without stating what it is compared to")
tf = box(s, 0.6, 1.65, 5.6, 4.9)
para(tf, "Foreground", size=19, bold=True, color=BLUE, first=True, space_after=8)
para(tf, f"The combined 2A set assembled earlier: Supplementary Table S2 plus the "
         f"paper's repository search, merged and de-duplicated — "
         f"{red['instances']:,} instances with a resolved position 2.",
     size=15, color=INK, space_after=10)
para(tf, "Databases: UniParc, UniProtKB, MGnify, IMG/VR.", size=13, color=INK2,
     space_after=16)
para(tf, "Background", size=19, bold=True, color=BLUE, space_after=8)
para(tf, f"{summ['background_proteins']:,} reviewed (Swiss-Prot) viral proteins, "
         f"{summ['background_residues']:,} residues. Viral, because 2A peptides are "
         f"overwhelmingly of viral origin.", size=15, color=INK, space_after=0)

x0 = 6.6
tf = box(s, x0, 1.65, 6.2, 0.5)
para(tf, "Three ways to ask 'what would you expect by chance?'", size=17, bold=True,
     color=INK, first=True, space_after=0)

for i, (name, colr, desc, val) in enumerate([
        ("1 · All amino acids, all positions", MUTED,
         "How common are D and E in viral protein sequence at all? The crudest "
         "expectation.", raw.loc[B_ALL, "background_DE_pct"]),
        ("2 · Position 2 of viral proteins", ORANGE,
         "The residue after the initiator Met. Controls for N-termini having their own "
         "composition, unlike the bulk average.", raw.loc[B_P2, "background_DE_pct"]),
        ("3 · Residue after any proline", GREEN,
         "P2 sits directly after a proline, and proline constrains what follows it. "
         "The best-matched control.", raw.loc[B_PP, "background_DE_pct"])]):
    y = 2.25 + i * 1.45
    panel(s, x0, y, 6.2, 1.28)
    tf = box(s, x0 + 0.2, y + 0.1, 4.5, 1.1)
    para(tf, name, size=14.5, bold=True, color=colr, first=True, space_after=4)
    para(tf, desc, size=12.5, color=INK2, space_after=0)
    tf = box(s, x0 + 4.75, y + 0.28, 1.3, 0.8)
    para(tf, f"{val:.1f}%", size=22, bold=True, color=colr, first=True, space_after=0,
         align=PP_ALIGN.CENTER)
footer(s, "Backgrounds computed from UniProt Swiss-Prot, taxonomy 10239 (Viruses), "
          "reviewed entries only")

# ============================================================ 4 the catch
s = slide()
header(s, "First, a counting problem that decides the answer",
       "The same number can read as 2.6% or 12.1% depending on what you count")
tf = box(s, 0.6, 1.65, 6.1, 4.6)
para(tf, "84% of the instances come from UniParc", size=18, bold=True, color=ORANGE,
     first=True, space_after=8)
para(tf, "UniParc archives every sequence variant separately, so a single protein can "
         "appear hundreds of times. Counting raw instances counts the same downstream "
         "context over and over.", size=15, color=INK, space_after=12)
para(tf, f"The {red['instances']:,} instances collapse to just "
         f"{red['unique_downstream_contexts']:,} distinct downstream contexts — a "
         f"{red['instances_per_unique_downstream_context']}-fold duplication.",
     size=15, color=INK, space_after=14)
para(tf, "Which unit is right?", size=18, bold=True, color=BLUE, space_after=8)
para(tf, "Position 2 is a property of the protein that follows the skip site, not of "
         "the 2A peptide. So each distinct downstream context should count once. "
         "That is the unit used for every result that follows.",
     size=15, color=INK, space_after=0)

stat(s, 7.0, 1.9, 2.7, "2.6%", "counting raw instances\n(redundancy included)",
     color=MUTED, vsize=32)
stat(s, 10.0, 1.9, 2.7, "12.1%", "counting distinct\ndownstream contexts",
     color=BLUE, vsize=32)
tf = box(s, 7.0, 3.7, 5.7, 2.6)
para(tf, "Why this matters", size=17, bold=True, color=INK, first=True, space_after=8)
para(tf, "The 2.4–2.6% figure reported previously is a database-redundancy artefact, "
         "not a biological depletion. It reflects how often a handful of well-deposited "
         "2A peptides were sequenced, not how often the motif occurs.",
     size=14.5, color=INK, space_after=8)
para(tf, "This corrects the earlier conclusion in 2a_combined_UBR3.", size=14.5,
     bold=True, color=AMBER, space_after=0)
footer(s, "Redundancy measured on the combined instance table; "
          "see enrichment_summary.json → redundancy")

# ============================================================ 5 the result
s = slide()
header(s, "Result", f"D/E at position 2, counting distinct downstream contexts "
                    f"(n = {int(prim.iloc[0]['foreground_n']):,})")
stat(s, 0.6, 1.6, 3.0, f"{prim.iloc[0]['observed_DE_pct']:.1f}%",
     "observed D/E at position 2\n95% CI "
     f"{prim.iloc[0]['observed_CI95_low_pct']:.1f}–{prim.iloc[0]['observed_CI95_high_pct']:.1f}%",
     color=BLUE)

rows = [("vs all amino acids, all positions", B_ALL, MUTED),
        ("vs residue after any proline", B_PP, GREEN),
        ("vs position 2 of viral proteins", B_P2, ORANGE)]
tf = box(s, 3.9, 1.55, 8.9, 0.4)
para(tf, "Compared with each background", size=15, bold=True, color=INK, first=True,
     space_after=0)
for i, (label, key, colr) in enumerate(rows):
    y = 1.95 + i * 1.15
    panel(s, 3.9, y, 8.9, 1.0)
    r = prim.loc[key]
    tf = box(s, 4.1, y + 0.08, 3.5, 0.9)
    para(tf, label, size=14, bold=True, color=colr, first=True, space_after=2)
    para(tf, f"expected {r['background_DE_pct']:.1f}%", size=12, color=INK2,
         space_after=0)
    tf = box(s, 7.7, y + 0.22, 1.5, 0.6)
    para(tf, f"{r['fold_change']:.2f}×", size=20, bold=True, color=colr, first=True,
         space_after=0, align=PP_ALIGN.CENTER)
    verdict = ("no difference" if r["fisher_p"] > 0.05
               else ("depleted" if r["fold_change"] < 1 else "enriched"))
    pv = "n.s." if r["fisher_p"] > 0.05 else f"p = {r['fisher_p']:.0e}"
    tf = box(s, 9.3, y + 0.22, 3.3, 0.6)
    para(tf, f"{verdict}   ({pv})", size=15, bold=True,
         color=(INK2 if r["fisher_p"] > 0.05 else colr), first=True, space_after=0)

panel(s, 0.6, 5.5, 12.2, 1.25)
tf = box(s, 0.85, 5.62, 11.7, 1.1)
para(tf, "In one sentence", size=14, bold=True, color=INK, first=True, space_after=4)
para(tf, "D/E follows the skip proline at essentially the rate chance predicts — there "
         "is no selection for or against the UBR3 motif at 2A skip sites. It is "
         "depleted only against viral N-termini, a background that is itself unusual.",
     size=15, color=INK, space_after=0)
footer(s, "Fisher exact test, two-sided, on 2×2 counts; "
          "exact binomial 95% CI on the observed proportion")

# ============================================================ 6 figure
s = slide()
header(s, "The analysis in one figure")
s.shapes.add_picture(str(HERE / "enrichment_figure.png"), Inches(0.6), Inches(1.55),
                     height=Inches(5.2))
for i, (lab, colr, title, body) in enumerate([
        ("A", BLUE, "Redundancy changes the answer",
         "Each bar is the same measurement counted a different way. As duplicate "
         "database entries are collapsed, D/E rises from 2.6% to 12.1% and crosses the "
         "two matched background lines."),
        ("B", INK, "Composition at position 2",
         "Blue is the observed 2A position 2; the three background sets sit beside it. "
         "D and E (shaded) track the grey and green backgrounds closely, and fall well "
         "below the orange viral-N-terminus background."),
        ("C", GREEN, "Fold change with 95% CI",
         "Two intervals cross 1.0 — no difference from chance. Only the viral "
         "position-2 comparison is significant, at 0.57×.")]):
    y = 1.62 + i * 1.72
    tf = box(s, 7.75, y, 5.1, 1.7)
    para(tf, f"{lab} · {title}", size=14.5, bold=True, color=colr, first=True,
         space_after=5)
    para(tf, body, size=12.5, color=INK, space_after=0)
footer(s, "Full-resolution vector version available as enrichment_figure.pdf / .svg")

# ============================================================ 7 interpretation
s = slide()
header(s, "What it means", "Three readings, in order of confidence")
items = [
    ("The motif occurs at chance frequency", BLUE,
     f"Against the two matched backgrounds — overall amino-acid composition "
     f"({prim.loc[B_ALL, 'background_DE_pct']:.1f}%) and the residue after any proline "
     f"({prim.loc[B_PP, 'background_DE_pct']:.1f}%) — the observed "
     f"{prim.iloc[0]['observed_DE_pct']:.1f}% is statistically indistinguishable "
     f"(1.10× and 1.09×, both n.s.). Nothing about the 2A skip favours or avoids "
     f"creating a UBR3 substrate."),
    ("The depletion vs viral N-termini is real but says less than it looks", ORANGE,
     f"Position 2 of viral proteins is {prim.loc[B_P2, 'background_DE_pct']:.1f}% D/E — "
     f"far above bulk composition. That is expected: Met-aminopeptidase only removes "
     f"the initiator methionine when residue 2 is small, so acidic residues accumulate "
     f"at position 2 of real N-termini. Comparing against it measures N-terminal "
     f"biology, not 2A biology."),
    ("The earlier 'uncommon motif' conclusion needs revising", AMBER,
     "The previously reported 2.4% came from raw instance counts dominated by UniParc "
     "duplication. On distinct downstream contexts the figure is 12.1% — five times "
     "higher, and exactly at background. The motif is not rare; it is ordinary."),
]
for i, (title, colr, body) in enumerate(items):
    y = 1.6 + i * 1.72
    panel(s, 0.6, y, 12.2, 1.55)
    num = s.shapes.add_shape(9, Inches(0.85), Inches(y + 0.35), Inches(0.55),
                             Inches(0.55))
    num.fill.solid()
    num.fill.fore_color.rgb = colr
    num.line.fill.background()
    num.shadow.inherit = False
    ntf = num.text_frame
    ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = ntf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(i + 1)
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT
    tf = box(s, 1.65, y + 0.14, 10.9, 1.35)
    para(tf, title, size=16.5, bold=True, color=colr, first=True, space_after=5)
    para(tf, body, size=13.5, color=INK, space_after=0)
footer(s, "Met-aminopeptidase specificity: initiator Met is cleaved only when residue 2 "
          "has a small side chain (A, C, G, P, S, T, V)")

# ============================================================ 8 caveats
s = slide()
header(s, "Caveats and what would strengthen this")
tf = box(s, 0.6, 1.7, 6.1, 4.9)
para(tf, "Limits of this analysis", size=19, bold=True, color=ORANGE, first=True,
     space_after=10)
for t in [
    "6.8% of instances still have no resolved downstream residue — almost all MGnify "
    "(MGYP), which needs the ~270 GB MGnify FASTA dump. resolve_mgnify.py is written "
    "but has not been run.",
    "The background is reviewed viral proteins. The 2A set also contains host and "
    "unclassified sequences; using a viral background for all of them is a "
    "simplification. The viral-only subset gives the same answer (12.5%).",
    "'Distinct downstream context' de-duplicates on the 20-aa product. Homologous "
    "proteins with slightly different downstream sequence still count separately, so "
    "some redundancy remains — this makes the estimate conservative, not inflated.",
]:
    para(tf, "•  " + t, size=13.5, color=INK, space_after=10)

tf = box(s, 7.0, 1.7, 5.8, 4.9)
para(tf, "What would make it stronger", size=19, bold=True, color=GREEN, first=True,
     space_after=10)
for t in [
    "Resolve the MGnify fraction and re-run — it is the one systematic gap.",
    "Test position 2 against a phylogeny-aware background rather than pooled "
    "composition, so related viruses do not count as independent observations.",
    "Ask the reverse question: among 2A sites that do carry P–[D/E], is anything else "
    "shared — host range, protein family, genome position?",
    "The 280 distinct P–[D/E] peptides already identified remain the useful output "
    "regardless of whether they are enriched: they are candidate UBR3 substrates.",
]:
    para(tf, "•  " + t, size=13.5, color=INK, space_after=10)
footer(s, "Scripts: enrichment.py, plot_enrichment.py, build_pptx.py · "
          "all outputs regenerate from these")

# ============================================================ 9 methods
s = slide()
header(s, "Methods", "So the numbers can be reproduced or challenged")
tf = box(s, 0.6, 1.65, 12.2, 5.0)
steps = [
    ("Foreground", "combined_all_instances.csv from 2a_combined_UBR3 — the merged "
     "Table S2 + repository set. Kept instances with position 2 resolved to one of the "
     "20 standard amino acids."),
    ("De-duplication", "Each distinct 20-aa downstream product counted once "
     f"({red['unique_downstream_contexts']:,} of {red['instances']:,} instances). "
     "Raw instances, unique accession and unique 2A core are also reported so the "
     "effect of the choice is visible."),
    ("Background", f"{summ['background_proteins']:,} reviewed viral proteins from "
     "UniProt (taxonomy 10239). Three counts: every residue; residue 2 of proteins "
     f"beginning with Met ({summ['background_proteins_starting_with_Met']:,}); and "
     "every residue immediately following a proline."),
    ("Statistics", "Two-sided Fisher exact test on the 2×2 table of D/E vs other, "
     "foreground vs background; chi-square reported alongside. Exact binomial 95% CI "
     "on the observed proportion. Fold change = observed % ÷ background %."),
]
first = True
for name, body in steps:
    p = para(tf, name, size=16, bold=True, color=BLUE, first=first, space_after=4)
    first = False
    para(tf, body, size=13.5, color=INK, space_after=12)
footer(s, "Outputs: enrichment_results.csv · enrichment_composition.csv · "
          "enrichment_summary.json · enrichment_figure.png/pdf/svg")

out = HERE / "2A_UBR3_enrichment.pptx"
prs.save(out)
print(f"-> {out.name}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")

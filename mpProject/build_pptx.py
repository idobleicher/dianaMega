"""
Build a PowerPoint deck summarising the MP-project enrichment analyses.

Reads the figures produced by:
  - mp_enrichment_heatmap.py   (vs MP-only library, results/vs_MP_library/)
  - mp_vs_full_library.py      (vs full library,    results/vs_full_library/)

Each slide carries:
  - a clear title
  - the figure
  - an "Explanation" paragraph (what the plot means)
  - a "Methods" paragraph (how it was computed)

Output: mpProject/MP_project_results.pptx
"""
from __future__ import annotations

from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
RES_MP = HERE / "results" / "vs_MP_library"
RES_FULL = HERE / "results" / "vs_full_library"
OUT = HERE / "MP_project_results.pptx"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x14, 0x2A, 0x4C)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF4, 0xF8)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title_band(slide, title: str, subtitle: str | None = None):
    """Coloured band at top of slide with title + optional subtitle."""
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.0))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.shadow.inherit = False

    tb = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.05), SLIDE_W - Inches(0.8), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    p.text = title
    p.runs[0].font.size = Pt(26)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if subtitle:
        sp = tf.add_paragraph()
        sp.text = subtitle
        sp.runs[0].font.size = Pt(14)
        sp.runs[0].font.color.rgb = RGBColor(0xCF, 0xDB, 0xE8)


def add_text_block(slide, left, top, width, height, header, body,
                   header_color=NAVY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = header
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(14)
    p.runs[0].font.color.rgb = header_color
    for line in body.split("\n"):
        bp = tf.add_paragraph()
        bp.text = line if line else " "
        bp.runs[0].font.size = Pt(11)
        bp.runs[0].font.color.rgb = GREY
        bp.space_before = Pt(2)
    return box


def add_image_fit(slide, image_path: Path, left, top, max_w, max_h):
    """Add image scaled to fit within (max_w, max_h) preserving aspect."""
    with Image.open(image_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    box_aspect = max_w / max_h
    if aspect > box_aspect:
        w = max_w
        h = int(max_w / aspect)
    else:
        h = max_h
        w = int(max_h * aspect)
    centred_left = left + (max_w - w) // 2
    centred_top = top + (max_h - h) // 2
    return slide.shapes.add_picture(
        str(image_path), centred_left, centred_top, width=w, height=h)


def add_figure_slide(prs, title: str, subtitle: str,
                     image_path: Path,
                     explanation: str, methods: str):
    """Layout: top band (title), left = image (~58% width), right = text (~42%)."""
    slide = add_blank_slide(prs)
    add_title_band(slide, title, subtitle)

    # background panel for the right text area
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(7.9), Inches(1.1),
        Inches(5.2), Inches(6.2))
    panel.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT

    add_image_fit(
        slide, image_path,
        left=Inches(0.25), top=Inches(1.2),
        max_w=Inches(7.5), max_h=Inches(6.0),
    )

    add_text_block(
        slide,
        left=Inches(8.05), top=Inches(1.3),
        width=Inches(4.95), height=Inches(2.7),
        header="Explanation", body=explanation,
    )
    add_text_block(
        slide,
        left=Inches(8.05), top=Inches(4.2),
        width=Inches(4.95), height=Inches(3.0),
        header="Methods", body=methods,
    )


def add_title_slide(prs, title: str, subtitle: str, body: str):
    slide = add_blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY

    tb = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.0), SLIDE_W - Inches(1.6), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.runs[0].font.size = Pt(44)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    sp = tf.add_paragraph()
    sp.text = subtitle
    sp.space_before = Pt(12)
    sp.runs[0].font.size = Pt(20)
    sp.runs[0].font.color.rgb = RGBColor(0xCF, 0xDB, 0xE8)

    for line in body.split("\n"):
        bp = tf.add_paragraph()
        bp.text = line if line else " "
        bp.space_before = Pt(8)
        bp.runs[0].font.size = Pt(14)
        bp.runs[0].font.color.rgb = RGBColor(0xE6, 0xEC, 0xF2)


def add_section_slide(prs, title: str, body_lines: list[str]):
    slide = add_blank_slide(prs)
    add_title_band(slide, title)

    tb = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.4), SLIDE_W - Inches(1.2),
        SLIDE_H - Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line if line.strip() else " "
        p.space_before = Pt(6)
        bullet = line.lstrip().startswith("• ")
        p.runs[0].font.size = Pt(15) if bullet else Pt(13)
        p.runs[0].font.color.rgb = GREY
        if line.strip().endswith(":"):
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = NAVY


# ---------------------------------------------------------------------------
# Slide content (titles + explanation + methods)
# ---------------------------------------------------------------------------
N_HITS = 155


def common_intro():
    return (
        "MP project = 155 user-supplied gene/transcript IDs.\n"
        "Library = N-terminome GPS screen library, union of sheets\n"
        "Data S3A / S3B / S3C, deduplicated on Ensembl Transcript ID\n"
        "(23,408 unique transcripts, 1,263 of which start with M-P)."
    )


def slide_specs_mp():
    """(title, subtitle, image, explanation, methods) tuples for vs-MP run."""
    return [
        ("Enrichment heatmap (primary)",
         "MP project (n=155) vs all MP peptides in the GPS library (n=1,263)",
         RES_MP / "mp_enrichment_heatmap.png",
         "Each cell shows log2(frequency in hits / frequency in MP "
         "library) for one amino acid (rows) at one residue position "
         "(columns).\n"
         "Red = enriched in your hits, blue = depleted, white ≈ no change.\n"
         "The strongest signals sit at positions 3–7: Y, M, F, W, V, C are "
         "enriched while K, R, N, P are depleted — the canonical "
         "hydrophobic / aromatic preference of an N-degron pathway after "
         "the conserved M-P start.",
         "Counts of each amino acid at each position were computed for both "
         "the hit set and the MP library, then converted to per-position "
         "frequencies. log2 ratio reported with a Laplace pseudocount of "
         "0.5 to avoid divide-by-zero. Cells where both groups have zero "
         "observations are forced to 0. Positions 1 (always M) and 2 "
         "(always P in MP peptides) are invariant in both groups and were "
         "excluded; positions 3–24 are shown."),

        ("Per-position frequency heatmaps",
         "Raw amino-acid frequencies — hits vs MP library",
         RES_MP / "mp_frequency_heatmaps.png",
         "Same data without the ratio — panel A is the hit set, panel B is "
         "the MP library. Both panels share the same colour scale so the "
         "darker cell is the more frequent residue at that position.\n"
         "The library (B) is roughly uniform; the hits (A) show clear "
         "patches of darker colour at hydrophobic / aromatic rows and at "
         "positions 3–7, which is what produces the red areas in the "
         "enrichment heatmap.",
         "Frequency = count of AA at position / total count at that "
         "position, expressed as a percentage. No pseudocount, no ratio "
         "— this is purely descriptive."),

        ("Statistical significance — Fisher's exact",
         "Signed -log10(p) for each (AA, position) cell",
         RES_MP / "mp_significance_heatmap.png",
         "Same layout as the enrichment heatmap, but each cell is a "
         "signed -log10(p) from a 2×2 Fisher's exact test (AA-of-interest "
         "vs others, hit set vs MP library).\n"
         "* = p<0.05, ** = p<0.01, *** = p<0.001 (two-sided).\n"
         "The most significant signals (***) are at positions 3 and 4: "
         "P and Y enriched, K / R / L / P depleted — these are the cells "
         "you can quote with statistical confidence.",
         "For each (AA, position) cell, a 2×2 contingency table was built "
         "(rows: hits / library, columns: this AA / all other AAs at that "
         "position). p-values come from scipy.stats.fisher_exact, "
         "two-sided. The colour is sign(observed − expected) × -log10(p). "
         "No multiple-testing correction is applied; treat single-cell "
         "p-values as descriptive."),

        ("Enrichment sequence logo",
         "Letter heights = log2(hit / library), coloured by chemistry",
         RES_MP / "mp_enrichment_logo.png",
         "A sequence logo where letter height encodes log2 enrichment "
         "directly. Letters above the baseline = enriched in hits; "
         "letters below (flipped upside-down) = depleted.\n"
         "Letters are coloured by amino-acid chemistry group: green = "
         "hydrophobic, purple = aromatic, light blue = polar, dark blue = "
         "basic, red = acidic, orange = special (C, G, P).\n"
         "The aromatic / hydrophobic stack at positions 3–7 above the "
         "baseline and the basic stack below are the visual summary of "
         "the substrate preference.",
         "Built with logomaker. Letter values are the same log2 enrichment "
         "matrix used in the heatmap (with the 0/0-cell mask applied), "
         "focused on positions 3–10 for figure readability. Negative values "
         "are rendered flipped below the baseline."),

        ("Property-grouped enrichment",
         "Amino acids collapsed into 6 chemistry groups",
         RES_MP / "mp_property_group_heatmap.png",
         "The same enrichment, summed within chemistry groups so the "
         "global trend is easier to read at a glance.\n"
         "Aromatic (FWY) is enriched at positions 3–10. Basic (KRH) is "
         "depleted across positions 3–10 — strongest at position 4 "
         "(log2 ≈ −2.5). Acidic (DE) is mostly flat. Special (CGP) shows "
         "the M-P motif extension into early positions.",
         "For each group, hit and library counts were summed across the "
         "AAs in the group at each position, then a log2 ratio was taken "
         "with the same pseudocount and total-normalisation as the AA-"
         "level heatmap. The groups are: Hydrophobic (AILMV), Aromatic "
         "(FWY), Polar (NQST), Basic (KRH), Acidic (DE), Special (CGP)."),

        ("Top movers per position (3–7)",
         "Four most enriched and four most depleted residues at each early position",
         RES_MP / "mp_top_movers.png",
         "For each early residue position (3 through 7), a horizontal "
         "bar chart of the four most enriched (top, positive bars) and "
         "four most depleted (bottom, negative bars) amino acids. Bar "
         "colour = chemistry group; numeric label = log2 enrichment.\n"
         "Useful as a 'pull-quote' figure: at every position 3–6 a "
         "basic residue (K or R) is the single most depleted residue, "
         "and at every position an aromatic / bulky-hydrophobic residue "
         "(Y, M, F, W, V, C) is the most enriched.",
         "Top-4 enriched and top-4 depleted residues at each position "
         "were ranked by the same log2-enrichment score. All five panels "
         "share the same x-axis range so the bar lengths are directly "
         "comparable across positions."),
    ]


def slide_specs_full():
    return [
        ("Enrichment heatmap (primary)",
         "MP project (n=155) vs full GPS library (n=23,408)",
         RES_FULL / "mp_enrichment_heatmap.png",
         "Same layout as the MP-vs-MP heatmap, but the background is now "
         "every N-terminus in the GPS library — not just the M-P starts.\n"
         "The biological signature is preserved: aromatic / bulky-"
         "hydrophobic residues (Y, M, F, W, V, C) are enriched at "
         "positions 3–7 and basic residues (K, R) are depleted, with the "
         "strongest depletion at position 4 (R log2 ≈ −3.0).\n"
         "Compared to the MP-vs-MP heatmap, the absolute log2 values are "
         "similar — confirming that the signal is biological and not an "
         "artefact of the library's M-P starting design.",
         "Counts → frequencies → log2 ratio with 0.5 pseudocount; cells "
         "with zero counts in both groups are masked to 0. Positions 1 "
         "(invariant M) and 2 (trivially enriched for P because every "
         "hit is an MP peptide while only ~5.4 % of the library is) are "
         "excluded. Positions 3–24 are shown."),

        ("Per-position frequency heatmaps",
         "Raw frequencies — hits vs full library",
         RES_FULL / "mp_frequency_heatmaps.png",
         "Side-by-side raw frequencies for hits (A) and the full library "
         "(B), shared colour scale. The library panel is more uniform "
         "and slightly more intense than the MP-only library because it "
         "contains ~19× more sequences.\n"
         "The hit panel still shows the aromatic / hydrophobic patches "
         "at positions 3–7 — the same signal that drives the enrichment "
         "heatmap.",
         "Frequency = AA count divided by total count at that position, "
         "as a percentage. No pseudocount, no ratio — purely descriptive."),

        ("Statistical significance — Fisher's exact",
         "Signed -log10(p), full-library background",
         RES_FULL / "mp_significance_heatmap.png",
         "Same significance map as the MP-vs-MP run, but with a larger "
         "library denominator (~23 k entries), so p-values are tighter "
         "and *** labels appear in more cells.\n"
         "The strongest *** signals reproduce the MP-vs-MP pattern: K, R "
         "depleted at positions 3–6; Y, P enriched at positions 3–4. "
         "Cells you can quote with statistical confidence (***) line up "
         "across both backgrounds — the biology is robust.",
         "Two-sided Fisher's exact for each (AA, position) cell against "
         "the full library. No multiple-testing correction. Positions 3–24."),

        ("Enrichment sequence logo",
         "Letter heights = log2(hit / full library)",
         RES_FULL / "mp_enrichment_logo.png",
         "Logo of the same enrichment matrix. The aromatic / "
         "hydrophobic stack above the baseline at positions 3–7 and the "
         "basic stack below are visually similar to the MP-vs-MP logo — "
         "just slightly compressed because the larger background "
         "smoothens single-residue tails.",
         "Same logomaker construction as the MP-vs-MP logo; positions "
         "3–10, 0/0-cell mask applied, negative letters flipped below "
         "the baseline."),

        ("Property-grouped enrichment",
         "Chemistry-group summary, full-library background",
         RES_FULL / "mp_property_group_heatmap.png",
         "Same chemistry-group summary as the MP-vs-MP run.\n"
         "Position 4 still shows the strongest depletion of Basic (KRH, "
         "log2 ≈ −2.4) and Aromatic (FWY) is enriched at positions 3–10. "
         "The fact that this pattern survives a 19× larger / more "
         "diverse background confirms it's a property of the hit set, "
         "not of the MP-only library.",
         "Same group definitions and computation as the MP-vs-MP run."),

        ("Top movers per position (3–7)",
         "Top-4 enriched / depleted residues, full-library background",
         RES_FULL / "mp_top_movers.png",
         "Per-position top movers, full-library background. Comparing "
         "side by side with the MP-vs-MP panel: the same residues "
         "occupy the top and bottom slots at almost every position — Y "
         "leads positions 3–5, K and R dominate the depleted side at "
         "positions 3–6, F and M follow Y in the enriched stack at "
         "positions 4–5.\n"
         "Quantitative differences exist (effects are typically a touch "
         "larger vs full library because the background is more "
         "diffuse), but the ranking is virtually identical.",
         "Same top-4 ranking and shared x-axis as the MP-vs-MP panel; "
         "positions 3–7."),
    ]


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Title slide
    add_title_slide(
        prs,
        title="MP Project — N-terminal Enrichment Analysis",
        subtitle="155 hit peptides vs N-terminome GPS screen library",
        body=(
            "Two-background comparison:\n"
            "  1) vs all MP-starting peptides in the library (n = 1,263)\n"
            "  2) vs the full GPS library, all N-termini (n = 23,408)\n"
            "\n"
            "Each section: enrichment heatmap, frequency heatmaps, "
            "Fisher's exact significance, sequence logo, "
            "property-grouped heatmap, per-position top movers."
        ),
    )

    # 2. Methods / data sources
    add_section_slide(
        prs,
        title="Data sources, definitions & methods",
        body_lines=[
            "Inputs:",
            "• MP project peptide list — 155 user-supplied gene / transcript "
            "IDs (some with explicit |ENST suffix). Loaded from "
            "peptide_list.txt; duplicate ' 2' suffixes are stripped.",
            "• Library — Data S3.xlsx (N-terminome GPS screen, in "
            "different genetic backgrounds). Sheets Data S3A / S3B / S3C "
            "are concatenated and deduplicated on Ensembl Transcript ID, "
            "giving 23,408 unique transcripts. 1,263 of those start "
            "with M-P. Each entry has a 24-residue N-terminal peptide "
            "sequence.",
            " ",
            "Definitions:",
            "• 'MP peptide' = peptide whose AA sequence starts with M-P.",
            "• Each user entry is matched to the library first by Ensembl "
            "transcript (when supplied), otherwise by gene symbol; we "
            "then take the first MP-starting peptide from the matching "
            "rows. All 155 entries match (0 unmatched).",
            " ",
            "Enrichment statistic:",
            "• At every (amino acid, position) cell, log2( hit_freq / "
            "library_freq ), with a Laplace pseudocount of 0.5 added to "
            "raw counts before normalisation.",
            "• Cells where both hit and library counts are exactly 0 are "
            "set to 0 (avoids pseudocount artefacts in fully empty cells).",
            "• Significance is two-sided Fisher's exact (AA vs others, "
            "hits vs library). No multiple-testing correction.",
        ],
    )

    # 3. vs MP library — section divider + 6 figure slides
    add_section_slide(
        prs,
        title="Section 1 — vs MP-only library background",
        body_lines=[
            "Background:",
            "• 1,263 unique MP-starting peptides in the GPS library.",
            "• This is the conservative comparison: it cancels out the "
            "library's design constraints around the M-P start and "
            "isolates the biological signal at positions 3 and beyond.",
            " ",
            "Why this matters:",
            "• If a residue is enriched at position 4 vs MP library, it "
            "is enriched relative to OTHER MP peptides — i.e. it is a "
            "true preference within the M-P starting population.",
        ],
    )
    for spec in slide_specs_mp():
        title, sub, img, expl, meth = spec
        add_figure_slide(prs, title, sub, img, expl, meth)

    # 4. vs full library — section divider + 6 figure slides
    add_section_slide(
        prs,
        title="Section 2 — vs FULL library background",
        body_lines=[
            "Background:",
            "• All 23,408 unique transcripts in the GPS library "
            "(every N-terminus, not just MP starts).",
            " ",
            "What this comparison answers:",
            "• Are the position-3-and-beyond preferences seen in "
            "Section 1 a property of the MP-only sub-library, or of the "
            "hit set? Comparing the same 155 hits to the entire library "
            "isolates the biology from the M-P start design.",
            " ",
            "Display choices (kept identical to Section 1):",
            "• Positions 3–24 are shown in every plot.",
            "• Position 1 is invariant (always M) in both groups.",
            "• Position 2 trivially enriches for P (every hit starts MP, "
            "only ~5.4 % of the library does); we exclude it so the "
            "biology at positions 3+ isn't obscured by the M-P design "
            "signal.",
        ],
    )
    for spec in slide_specs_full():
        title, sub, img, expl, meth = spec
        add_figure_slide(prs, title, sub, img, expl, meth)

    # 5. Final summary slide
    add_section_slide(
        prs,
        title="Summary of findings",
        body_lines=[
            "Background-invariant biological signal (positions 3–7):",
            "• Strong enrichment of aromatic / bulky-hydrophobic residues "
            "(Y, M, L, F, W, V, C) at the early peptide positions.",
            "• Strong depletion of basic residues (K, R, H) and several "
            "polar / amide residues (N, Q) at the same positions.",
            "• Largest single signals: K depleted at positions 3–6 "
            "(log2 ≈ −2 to −2.6); R depleted at position 4 (≈ −3.0 vs "
            "full library, −3.1 vs MP library).",
            " ",
            "Robustness across backgrounds:",
            "• The same residues occupy the same enriched / depleted "
            "slots whether the hits are compared against the MP-only "
            "library (n = 1,263) or the full library (n = 23,408).",
            "• Effect magnitudes are similar between the two backgrounds; "
            "the full-library comparison gives slightly tighter p-values "
            "thanks to the larger denominator.",
            " ",
            "Outputs delivered:",
            "• mpProject/results/vs_MP_library/   — 6 figures + CSVs",
            "• mpProject/results/vs_full_library/ — 6 figures + CSVs",
            "• MP_project_results.pptx (this deck)",
        ],
    )

    prs.save(OUT)
    print(f"Saved {OUT}")


if __name__ == "__main__":
    build()
